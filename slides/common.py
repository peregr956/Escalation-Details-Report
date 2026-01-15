"""Common slide layouts and utilities.

This module contains shared slide creation functions used across sections:
- Title slide layout (gradient background, centered text)
- Section header layout (gray background, narrative subtitle)
- Content slide layout (standard header/footer with content area)
- Chart data preparation (format data for Chart.js rendering)
- Chart image insertion (place rendered chart PNGs into slides)

REFACTORING STATUS:
    Most implementations remain in generate_presentation.py.
    This module provides re-exports for gradual migration.

    Step 1 (current): Re-export from generate_presentation.py
    Step 2 (next): Move implementations here, update generate_presentation.py
                   to import from this module instead
    Step 3 (final): Remove re-exports, direct imports only

Usage:
    from slides.common import create_section_header_layout
    # or
    from generate_presentation import create_section_header_layout  # until migrated
"""

# Re-export from generate_presentation for now
# When implementations are moved here, update these imports
def _lazy_import_from_main():
    """Lazy import to avoid circular imports during transition."""
    import sys
    from pathlib import Path

    # Add parent directory to path
    parent = Path(__file__).parent.parent
    if str(parent) not in sys.path:
        sys.path.insert(0, str(parent))

    from generate_presentation import (
        create_section_header_layout as _create_section_header_layout,
        prepare_chart_data as _prepare_chart_data,
    )
    return {
        'create_section_header_layout': _create_section_header_layout,
        'prepare_chart_data': _prepare_chart_data,
    }

# These will be populated on first use
_imported_functions = {}


def create_section_header_layout(prs, section_title, narrative_subtitle=None):
    """Create a section header slide (gray card between sections).

    Args:
        prs: Presentation object
        section_title: Main title text (e.g., "Executive Summary")
        narrative_subtitle: Narrative text below title

    Returns:
        The created slide
    """
    if 'create_section_header_layout' not in _imported_functions:
        _imported_functions.update(_lazy_import_from_main())
    return _imported_functions['create_section_header_layout'](
        prs, section_title, narrative_subtitle
    )


def prepare_chart_data(data):
    """Prepare chart data from ReportData for Chart.js rendering.

    Args:
        data: ReportData instance

    Returns:
        Dict with chart-specific data structures
    """
    if 'prepare_chart_data' not in _imported_functions:
        _imported_functions.update(_lazy_import_from_main())
    return _imported_functions['prepare_chart_data'](data)


# Functions that will be implemented directly (not re-exports)
def create_title_slide_layout(prs, title_text, subtitle_text=None):
    """Create a title slide with gradient background.

    NOTE: Implementation pending migration from generate_presentation.py

    Args:
        prs: Presentation object
        title_text: Main title
        subtitle_text: Optional subtitle

    Returns:
        The created slide
    """
    raise NotImplementedError(
        "create_title_slide_layout not yet migrated. "
        "Use generate_presentation.build_executive_summary_slides() instead."
    )


def create_content_slide_layout(prs, slide_title, content_items=None):
    """Create a standard content slide with header/footer.

    NOTE: Implementation pending migration from helpers.py

    Args:
        prs: Presentation object
        slide_title: Slide title text
        content_items: Optional list of content items

    Returns:
        Tuple of (slide, content_top)
    """
    # This one we can implement by delegating to helpers
    from helpers import setup_content_slide
    return setup_content_slide(prs, slide_title)


def insert_chart_image(slide, image_path, left, top, width=None, height=None):
    """Insert a rendered chart image into a slide.

    Args:
        slide: Slide object
        image_path: Path to the PNG image
        left: Left position (Inches)
        top: Top position (Inches)
        width: Optional width (Inches)
        height: Optional height (Inches)

    Returns:
        True if successful, False otherwise
    """
    from pathlib import Path
    from pptx.util import Inches
    import logging

    logger = logging.getLogger(__name__)
    image_path = Path(image_path)

    if not image_path.exists():
        logger.warning(f"Chart image not found: {image_path}")
        return False

    try:
        if width and height:
            slide.shapes.add_picture(str(image_path), left, top, width, height)
        elif width:
            slide.shapes.add_picture(str(image_path), left, top, width=width)
        elif height:
            slide.shapes.add_picture(str(image_path), left, top, height=height)
        else:
            slide.shapes.add_picture(str(image_path), left, top)
        return True
    except Exception as e:
        logger.error(f"Failed to insert chart image: {e}")
        return False


__all__ = [
    'create_title_slide_layout',
    'create_section_header_layout',
    'create_content_slide_layout',
    'prepare_chart_data',
    'insert_chart_image',
]
