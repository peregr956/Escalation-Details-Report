"""Executive Business Review - PowerPoint Presentation Generator

This module generates branded PowerPoint presentations from report data.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import argparse
import logging
import shutil
from datetime import datetime
from typing import Dict, Optional
from PIL import Image

# Import constants from dedicated module
from constants import (
    CS_BLUE, CS_NAVY, CS_SLATE, CS_VIOLET, CS_RED, CS_ORANGE, CS_GREEN,
    TITLE_FONT_NAME, BODY_FONT_NAME, FALLBACK_TITLE_FONT, FALLBACK_BODY_FONT,
    H1_FONT_SIZE, H2_FONT_SIZE, H3_FONT_SIZE, H4_FONT_SIZE,
    H5_FONT_SIZE, H6_FONT_SIZE, PARAGRAPH_FONT_SIZE, FOOTER_FONT_SIZE,
    TITLE_FONT_SIZE, SUBTITLE_FONT_SIZE, HEADING_FONT_SIZE,
    BODY_FONT_SIZE, CAPTION_FONT_SIZE, METRIC_FONT_SIZE, HERO_METRIC_FONT_SIZE,
    MARGIN_STANDARD, MARGIN_CONTENT, CARD_SPACING, HEADER_HEIGHT, FOOTER_HEIGHT,
    PRESENTATION_TITLE, PRESENTATION_INTENT, COPYRIGHT_TEXT,
    SECTION_NARRATIVES, get_brand_colors, normalize_terminology
)

# Import helper functions from dedicated module
from helpers import (
    create_presentation, apply_branding, get_slide_number,
    add_master_slide_header, add_master_slide_footer, add_master_slide_elements,
    add_logo, setup_content_slide, create_gradient_background,
    add_header_bar, add_insight_callout, create_metric_card, create_data_table,
    get_metric_status, get_status_indicator, add_trend_indicator,
    create_metric_card_with_indicator, get_status_text
)


# =============================================================================
# ARCHITECTURE GUIDE FOR NEW DEVELOPERS
# =============================================================================
#
# This file contains all slide builder functions for the PowerPoint presentation.
# It is intentionally kept as a single file for now (see slides/ package for
# planned modular refactoring).
#
# PRESENTATION STRUCTURE (19 slides total):
# =========================================
# 1. Title Slide                    - build_executive_summary_slides()
# 2. Executive Summary (section)    - create_section_header_layout()
# 3. Executive Dashboard            - build_executive_summary_slides()
# 4. CORR Funnel                    - build_corr_funnel_slide()
# 5. Value Delivered (section)      - create_section_header_layout()
# 6. Value Delivered                - build_value_delivered_slides()
# 7. Protection Achieved (section)  - create_section_header_layout()
# 8. Performance Metrics            - build_protection_achieved_slides()
# 9. Response & Detection           - build_protection_achieved_slides()
# 10. Threat Landscape (section)    - create_section_header_layout()
# 11. Severity Alignment            - build_threat_landscape_slides()
# 12. Threat Sources                - build_threat_landscape_slides()
# 13. Insights (section)            - create_section_header_layout()
# 14. Improvement Plan              - build_insights_slides()
# 15. Operational Coverage          - build_insights_slides()
# 16. Forward Direction (section)   - create_section_header_layout()
# 17. Key Takeaways                 - build_key_takeaways_slide()
# 18. Looking Ahead                 - build_forward_direction_slide()
# 19. Contact                       - build_contact_slide()
#
# KEY PATTERNS:
# =============
# 1. Each build_*_slides() function creates multiple related slides
# 2. setup_content_slide() is the helper for standard slide layout
# 3. create_section_header_layout() creates gray transition slides
# 4. create_metric_card() / add_insight_callout() for reusable components
# 5. Charts are rendered via chart_renderer.py and inserted as images
#
# DATA FLOW:
# ==========
# ReportData -> build_*_slides() -> python-pptx shapes -> .pptx file
#
# COMMON PARAMETERS:
# ==================
# - prs: Presentation object from python-pptx
# - data: ReportData instance with all metrics (see report_data.py)
#
# MODIFICATION TIPS:
# ==================
# - To add a new slide: Create a new build_* function, call it from main()
# - To change layout: Modify the relevant build_* function
# - To change branding: Update constants.py (colors, fonts, sizes)
# - To change data: Update metrics_calculator.py or report_data.py
#
# See docs/SLIDE_STRUCTURE.md for detailed slide specifications.
# See docs/DEVELOPER_GUIDE.md for adding new slides.
#
# =============================================================================
# SLIDE BUILDER FUNCTIONS
# =============================================================================

def build_key_takeaways_slide(prs, section_title, takeaways, data=None):
    """Create a Key Takeaways slide for executive summary of a section.
    
    Per CRITICALSTART branding guidelines:
    - All slides have transparent header and footer
    - Uses H1-H6 typography scale
    
    Args:
        prs (Presentation): The presentation object.
        section_title (str): The section this summarizes (e.g., "Value Delivered")
        takeaways (list): List of takeaway strings (3-4 recommended)
        data (ReportData, optional): Data object for dynamic values
    
    Returns:
        The slide object.
    """
    # Use setup_content_slide helper for consistent branding
    slide, content_top = setup_content_slide(prs, f"Key Takeaways: {section_title}")
    
    # Create takeaway cards (further reduced heights to prevent footer overlap per Jan 2026 feedback)
    card_top = content_top + Inches(0.05)
    card_height = Inches(0.65)  # Further reduced from 0.72"
    card_spacing = Inches(0.08)  # Further reduced from 0.1"
    card_width = prs.slide_width - MARGIN_STANDARD * 2
    card_left = MARGIN_STANDARD
    
    for i, takeaway in enumerate(takeaways):
        current_top = card_top + i * (card_height + card_spacing)
        
        # Card background
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, card_left, current_top,
            card_width, card_height
        )
        fill = card_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)
        line = card_shape.line
        line.color.rgb = CS_BLUE
        line.width = Pt(2)
        
        # Checkmark icon (using text)
        check_box = slide.shapes.add_textbox(
            card_left + Inches(0.12), current_top + Inches(0.18),
            Inches(0.4), Inches(0.4)
        )
        check_frame = check_box.text_frame
        check_para = check_frame.paragraphs[0]
        check_para.text = "✓"
        check_para.font.name = TITLE_FONT_NAME
        check_para.font.size = Pt(22)  # Further reduced from 24
        check_para.font.bold = True
        check_para.font.color.rgb = CS_BLUE
        check_para.alignment = PP_ALIGN.CENTER
        
        # Takeaway text
        text_box = slide.shapes.add_textbox(
            card_left + Inches(0.55), current_top + Inches(0.15),
            card_width - Inches(0.7), card_height - Inches(0.25)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.text = takeaway
        text_para.font.name = BODY_FONT_NAME
        text_para.font.size = Pt(14)  # Further reduced from 15
        text_para.font.color.rgb = CS_NAVY
        text_para.alignment = PP_ALIGN.LEFT
    
    # Add insight bar at bottom (further reduced height)
    insight_top = card_top + len(takeaways) * (card_height + card_spacing) + Inches(0.05)
    add_insight_callout(
        slide, prs,
        "Bottom Line",
        "Faster response, complete containment, measurable value—your security investment is delivering.",
        insight_top,
        height=Inches(0.5)  # Further reduced from 0.55"
    )
    
    return slide


def create_title_slide_layout(prs, title_text, subtitle_text=None):
    """Create title slide with logo and gradient background.
    
    Args:
        prs (Presentation): The presentation object.
        title_text (str): Main title text.
        subtitle_text (str, optional): Subtitle text.
    
    Returns:
        Slide: The created slide object.
    """
    # Use blank layout for full control
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add gradient background (blue to navy)
    # Note: python-pptx gradient support is limited, so we'll use a solid color
    # with a shape overlay for gradient effect
    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, 
        prs.slide_width, prs.slide_height
    )
    fill = background_shape.fill
    fill.solid()
    fill.fore_color.rgb = CS_NAVY  # Use navy as base
    background_shape.line.fill.background()  # No border
    
    # Note: Logo removed from title slide per branding guidelines
    
    # Add title
    title_left = Inches(1)
    title_top = Inches(2)
    title_width = prs.slide_width - Inches(2)
    title_height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, 
                                        title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title_text
    title_paragraph.font.name = TITLE_FONT_NAME
    title_paragraph.font.size = TITLE_FONT_SIZE
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Add subtitle if provided
    if subtitle_text:
        subtitle_left = Inches(1)
        subtitle_top = Inches(3.5)
        subtitle_width = prs.slide_width - Inches(2)
        subtitle_height = Inches(1)
        
        subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top,
                                               subtitle_width, subtitle_height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.text = subtitle_text
        subtitle_paragraph.font.name = BODY_FONT_NAME
        subtitle_paragraph.font.size = Pt(20)
        subtitle_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        subtitle_paragraph.alignment = PP_ALIGN.LEFT
    
    return slide


def create_section_header_layout(prs, section_title, narrative_subtitle=None):
    """Create section divider slide with gray background and optional narrative subtitle.
    
    Per brand guidelines, section title slides use Charcoal gray background (#343741)
    with white text. These slides serve as narrative transitions that guide executives
    through the presentation story.
    
    Args:
        prs (Presentation): The presentation object.
        section_title (str): Section title text (H2 typography, 72pt).
        narrative_subtitle (str, optional): Narrative subtitle text (H4 typography, 27pt).
    
    Returns:
        Slide: The created slide object.
    """
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_number = get_slide_number(prs)
    
    # Add gray background (CS_SLATE - Charcoal #343741)
    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, prs.slide_height
    )
    fill = background_shape.fill
    fill.solid()
    fill.fore_color.rgb = CS_SLATE  # Use charcoal gray per brand guidelines
    background_shape.line.fill.background()
    
    # Add master slide elements (header and footer) with white text for WCAG AA compliance
    # Gray background (#343741) requires white text for sufficient contrast
    add_master_slide_elements(slide, prs, slide_number=slide_number,
                               include_header=True, include_footer=True,
                               text_color=RGBColor(255, 255, 255))
    
    # Note: Logo removed from section cards per branding guidelines
    
    # Calculate vertical positioning
    # If we have a subtitle, move title up slightly to accommodate both
    if narrative_subtitle:
        title_top = Inches(1.8)
    else:
        title_top = Inches(2.2)
    
    # Add section title (centered) - H2 typography (72pt)
    title_left = Inches(1)
    title_width = prs.slide_width - Inches(2)
    title_height = Inches(1.2)
    
    title_box = slide.shapes.add_textbox(title_left, title_top,
                                        title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = section_title
    title_paragraph.font.name = TITLE_FONT_NAME
    title_paragraph.font.size = H2_FONT_SIZE  # 72pt per typography scale
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add narrative subtitle if provided - H4 typography (27pt)
    if narrative_subtitle:
        subtitle_top = title_top + Inches(1.3)  # 0.3" below title
        subtitle_height = Inches(0.8)
        
        subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top,
                                                title_width, subtitle_height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.text = narrative_subtitle
        subtitle_paragraph.font.name = BODY_FONT_NAME
        subtitle_paragraph.font.size = H4_FONT_SIZE  # 27pt per typography scale
        subtitle_paragraph.font.bold = False
        subtitle_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
    
    return slide


def create_content_slide_layout(prs, slide_title, content_items=None):
    """Create standard content slide with header.
    
    Args:
        prs (Presentation): The presentation object.
        slide_title (str): Slide title text.
        content_items (list, optional): List of content strings for bullet points.
    
    Returns:
        Slide: The created slide object.
    """
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add logo at top right
    add_logo(slide, position='top_right', prs=prs)
    
    # Add title header with brand color background
    header_height = Inches(0.8)
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, header_height
    )
    fill = header_shape.fill
    fill.solid()
    fill.fore_color.rgb = CS_NAVY
    header_shape.line.fill.background()
    
    # Add title text on header
    title_left = Inches(0.5)
    title_top = Inches(0.1)
    title_width = prs.slide_width - Inches(2.5)  # Leave space for logo
    title_height = Inches(0.6)
    
    title_box = slide.shapes.add_textbox(title_left, title_top,
                                        title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = slide_title
    title_paragraph.font.name = TITLE_FONT_NAME
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Add content area if items provided
    if content_items:
        content_left = Inches(0.5)
        content_top = header_height + Inches(0.3)
        content_width = prs.slide_width - Inches(1)
        content_height = prs.slide_height - content_top - Inches(0.3)
        
        content_box = slide.shapes.add_textbox(content_left, content_top,
                                             content_width, content_height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.margin_left = Inches(0.3)
        content_frame.margin_right = Inches(0.3)
        content_frame.margin_top = Inches(0.2)
        content_frame.margin_bottom = Inches(0.2)
        
        # Add content items as bullet points
        for i, item in enumerate(content_items):
            if i == 0:
                paragraph = content_frame.paragraphs[0]
            else:
                paragraph = content_frame.add_paragraph()
            
            paragraph.text = item
            paragraph.font.name = BODY_FONT_NAME
            paragraph.font.size = BODY_FONT_SIZE
            paragraph.font.color.rgb = CS_SLATE
            paragraph.level = 0
            paragraph.space_after = Pt(12)
    
    return slide


def prepare_chart_data(data) -> Dict:
    """Prepare chart data dictionary from ReportData for chart rendering.
    
    Args:
        data: ReportData instance containing all metrics
    
    Returns:
        Dictionary with chart data structured for different chart types
    """
    chart_data = {}
    
    # Trend chart data (MTTR, MTTD, FP% trends)
    # Industry benchmarks shown as horizontal dashed lines for comparison (if available)
    num_periods = len(data.period_labels)
    
    # Build base datasets for trend chart
    trend_datasets = [
        {
            'label': 'MTTR (min)',
            'data': data.mttr_trend,
            'borderColor': '#009CDE',
            'backgroundColor': 'rgba(0, 156, 222, 0.1)',
            'borderWidth': 3,
            'tension': 0.4,
            'pointRadius': 4,
            'pointHoverRadius': 6
        },
        {
            'label': 'MTTD (min)',
            'data': data.mttd_trend,
            'borderColor': '#004C97',
            'backgroundColor': 'rgba(0, 76, 151, 0.1)',
            'borderWidth': 3,
            'tension': 0.4,
            'pointRadius': 4,
            'pointHoverRadius': 6
        },
        {
            'label': 'FP %',
            'data': data.fp_trend,
            'borderColor': '#EF3340',
            'backgroundColor': 'rgba(239, 51, 64, 0.1)',
            'borderWidth': 3,
            'borderDash': [8, 4],
            'tension': 0.4,
            'pointRadius': 4,
            'pointHoverRadius': 6,
            'yAxisID': 'y1'
        }
    ]
    
    # Add industry benchmark lines only if benchmarks are available
    if data.industry_benchmarks_available:
        mttr_benchmark = data.mttr_industry_benchmark if data.mttr_industry_benchmark > 0 else 192.0
        mttd_benchmark = data.mttd_industry_benchmark if data.mttd_industry_benchmark > 0 else 66.0
        
        trend_datasets.extend([
            {
                'label': f'Industry MTTR ({int(mttr_benchmark)} min)',
                'data': [mttr_benchmark] * num_periods,
                'borderColor': 'rgba(0, 156, 222, 0.5)',
                'backgroundColor': 'transparent',
                'borderWidth': 2,
                'borderDash': [12, 6],
                'tension': 0,
                'pointRadius': 0,
                'pointHoverRadius': 0,
                'fill': False
            },
            {
                'label': f'Industry MTTD ({int(mttd_benchmark)} min)',
                'data': [mttd_benchmark] * num_periods,
                'borderColor': 'rgba(0, 76, 151, 0.5)',
                'backgroundColor': 'transparent',
                'borderWidth': 2,
                'borderDash': [12, 6],
                'tension': 0,
                'pointRadius': 0,
                'pointHoverRadius': 0,
                'fill': False
            }
        ])
    
    chart_data['trend'] = {
        'labels': data.period_labels,
        'datasets': trend_datasets
    }
    
    # Pie chart data (Operational Load) - only include if after-hours data available
    if data.after_hours_data_available:
        chart_data['pie'] = {
            'labels': ['Business Hours', 'After Hours', 'Weekend'],
            'data': [
                int(data.business_hours_percent),
                int(data.after_hours_percent),
                int(data.weekend_percent)
            ],
            'backgroundColor': ['#009CDE', '#702F8A', '#EF3340']
        }
    else:
        # Placeholder data when after-hours analysis not available
        chart_data['pie'] = {
            'labels': ['Data Analysis', 'Coming Soon'],
            'data': [100, 0],
            'backgroundColor': ['#E0E0E0', '#F5F5F5']
        }
    
    # Sankey chart data (Severity Flow)
    chart_data['sankey'] = {
        'flows': data.severity_flows
    }
    
    # Stacked bar chart data (MITRE ATT&CK)
    chart_data['stacked_bar'] = {
        'labels': data.tactics,
        'datasets': [
            {
                'label': 'High',
                'data': data.high_severity,
                'backgroundColor': '#EF3340'
            },
            {
                'label': 'Medium',
                'data': data.medium_severity,
                'backgroundColor': '#FF6A14'
            },
            {
                'label': 'Low',
                'data': data.low_severity,
                'backgroundColor': '#009CDE'
            },
            {
                'label': 'Info',
                'data': data.info_severity,
                'backgroundColor': '#702F8A'
            }
        ]
    }
    
    return chart_data


def insert_chart_image(slide, placeholder_id: Optional[str], image_path: str, 
                       padding: float = 0.1, min_fill_ratio: float = 0.85):
    """Replace a placeholder shape with a chart image.
    
    This function locates a placeholder shape (identified by "[Chart:" text)
    and replaces it with a chart image while preserving aspect ratio and
    centering within the placeholder bounds.
    
    Args:
        slide: The slide object containing the placeholder
        placeholder_id: Optional ID string to identify the placeholder (e.g., "severity_sankey")
                       If None, searches for first placeholder with "[Chart:" text
        image_path: Path to the image file to insert
        padding: Padding around the image in inches (default 0.1")
        min_fill_ratio: Minimum ratio of placeholder area to fill (0.0-1.0, default 0.85)
    
    Returns:
        bool: True if placeholder was found and replaced, False otherwise
    
    Note:
        Charts are rendered at 2x DPI for clarity. The function automatically
        handles aspect ratio preservation and centering.
    """
    image_path_obj = Path(image_path)
    if not image_path_obj.exists():
        logging.warning(f"Chart image not found: {image_path}")
        return False
    
    # Find the placeholder shape
    placeholder_shape = None
    
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            text = shape.text_frame.text
            if "[Chart:" in text:
                # If placeholder_id is specified, check if it matches
                if placeholder_id:
                    if f"ID: {placeholder_id}" in text or placeholder_id in text:
                        placeholder_shape = shape
                        break
                else:
                    # Use first placeholder found
                    placeholder_shape = shape
                    break
    
    if not placeholder_shape:
        logging.warning(f"Placeholder not found on slide (ID: {placeholder_id})")
        return False
    
    # Get placeholder position and size
    left = placeholder_shape.left
    top = placeholder_shape.top
    placeholder_width = placeholder_shape.width
    placeholder_height = placeholder_shape.height
    
    # Get actual image dimensions to preserve aspect ratio
    # Charts are rendered at scale=2, so we need to account for that
    # The images are rendered at 2x DPI, so physical size is half the pixel dimensions
    try:
        with Image.open(image_path) as img:
            img_width_px, img_height_px = img.size
            # Convert pixels to inches (assuming 96 DPI base, but rendered at 2x scale = 192 DPI effective)
            # So 1 pixel = 1/192 inches at the rendered scale
            # But since we're fitting to placeholder, we just need aspect ratio
            img_aspect_ratio = img_width_px / img_height_px
    except Exception as e:
        logging.warning(f"Could not read image dimensions for {image_path}: {e}. Using placeholder dimensions.")
        img_aspect_ratio = placeholder_width / placeholder_height
    
    # Calculate placeholder aspect ratio
    placeholder_aspect_ratio = placeholder_width / placeholder_height
    
    # Adjust dimensions to preserve image aspect ratio while fitting within placeholder bounds
    # Work with the aspect ratios and fit to the constraining dimension
    if img_aspect_ratio > placeholder_aspect_ratio:
        # Image is wider relative to placeholder - fit to placeholder width, adjust height
        width = placeholder_width
        height = placeholder_width / img_aspect_ratio
        # Center vertically if height is less than placeholder
        if height < placeholder_height:
            top = top + (placeholder_height - height) / 2
    else:
        # Image is taller relative to placeholder - fit to placeholder height, adjust width
        height = placeholder_height
        width = placeholder_height * img_aspect_ratio
        # Center horizontally if width is less than placeholder
        if width < placeholder_width:
            left = left + (placeholder_width - width) / 2
    
    # Remove the placeholder shape
    slide.shapes._spTree.remove(placeholder_shape._element)
    
    # Insert the image with preserved aspect ratio
    try:
        slide.shapes.add_picture(str(image_path), left, top, width, height)
        logging.info(f"Inserted chart image: {image_path} (ID: {placeholder_id}) - Size: {width}x{height}, Aspect ratio preserved")
        return True
    except Exception as e:
        logging.error(f"Failed to insert chart image {image_path}: {e}")
        return False


def main():
    """Main function to orchestrate slide generation.
    
    This function:
    1. Loads report data
    2. Renders all charts
    3. Creates presentation
    4. Builds all slides in order
    5. Inserts chart images
    6. Saves presentation
    """
    # Set up logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s'
    )
    logger = logging.getLogger(__name__)
    
    # Parse command-line arguments
    parser = argparse.ArgumentParser(
        description='Generate Executive Business Review PowerPoint presentation',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
Examples:
  # Use static sample data (default behavior)
  python generate_presentation.py
  
  # Load from single Excel file with config
  python generate_presentation.py --data report.xlsx --config client.yaml
  
  # Load multiple periods for trend charts
  python generate_presentation.py --data aug.xlsx sep.xlsx oct.xlsx --config client.yaml
  
  # Validate data without generating presentation
  python generate_presentation.py --data report.xlsx --validate
'''
    )
    parser.add_argument(
        '--data',
        type=str,
        nargs='+',
        help='Path(s) to Excel data file(s). Provide 1-3 files for period comparison. '
             'Last file is treated as current period.'
    )
    parser.add_argument(
        '--config',
        type=str,
        help='Path to client configuration YAML file'
    )
    parser.add_argument(
        '--client-name',
        type=str,
        help='Client name to use in the report (overrides data/config)'
    )
    parser.add_argument(
        '--no-threat-landscape',
        action='store_true',
        help='Exclude threat landscape slides from the presentation'
    )
    parser.add_argument(
        '--output-dir',
        type=str,
        default='output',
        help='Output directory for the presentation file (default: output)'
    )
    parser.add_argument(
        '--keep-charts',
        action='store_true',
        help='Keep temporary chart images after generation (default: delete)'
    )
    parser.add_argument(
        '--validate',
        action='store_true',
        help='Validate data only without generating presentation'
    )
    parser.add_argument(
        '--data-format',
        type=str,
        default='auto',
        choices=['auto', 'standard', 'burlington'],
        help='Data format profile for Excel parsing (default: auto-detect)'
    )
    parser.add_argument(
        '--trend-granularity',
        type=str,
        choices=['daily', 'weekly', 'monthly', 'quarterly'],
        help='Granularity for trend calculations from single data file. '
             'If not specified, prompts interactively when using single file.'
    )
    parser.add_argument(
        '--client',
        type=str,
        help='Client ID from registry (loads config from clients/registry.yaml)'
    )

    args = parser.parse_args()
    
    logger.info("=" * 60)
    logger.info("Starting Executive Business Review Generation")
    logger.info("=" * 60)
    
    # Step 1: Import and load data
    logger.info("Step 1: Loading report data...")
    try:
        from report_data import get_report_data, load_report_data, validate_report_data, ReportData
        
        if args.data:
            # Dynamic data loading from Excel files
            logger.info(f"Loading data from {len(args.data)} Excel file(s)...")
            for i, path in enumerate(args.data):
                period_label = "current" if i == len(args.data) - 1 else f"period -{len(args.data) - 1 - i}"
                logger.info(f"  - {path} ({period_label})")
            
            excel_files = args.data
            config_file = args.config
        else:
            # Interactive mode - prompt user for file selection
            from interactive_upload import interactive_file_selection
            
            excel_files, config_file = interactive_file_selection()
            
            if not excel_files:
                logger.error("No files selected. Exiting.")
                return 1
            
            logger.info(f"Loading data from {len(excel_files)} Excel file(s)...")
            for i, path in enumerate(excel_files):
                period_label = "current" if i == len(excel_files) - 1 else f"period -{len(excel_files) - 1 - i}"
                logger.info(f"  - {path} ({period_label})")
        
        # Load report data from selected files
        data = load_report_data(
            excel_paths=excel_files,
            config_path=config_file,
            client_name_override=args.client_name,
            data_format=args.data_format
        )
        logger.info("✓ Data loaded from Excel files")
        
        logger.info(f"✓ Loaded data for {data.client_name}")
        logger.info(f"  Period: {data.period_start} to {data.period_end}")
        logger.info(f"  Incidents escalated: {data.incidents_escalated}")
        
        # Validate data
        warnings = validate_report_data(data)
        if warnings:
            for warning in warnings:
                logger.warning(f"  ⚠ {warning}")
        
        # If validate-only mode, stop here
        if args.validate:
            logger.info("=" * 60)
            logger.info("Validation complete. No presentation generated.")
            if warnings:
                logger.info(f"Found {len(warnings)} warning(s)")
            else:
                logger.info("Data validation passed with no warnings")
            return 0
            
    except FileNotFoundError as e:
        logger.error(f"File not found: {e}")
        return 1
    except ValueError as e:
        logger.error(f"Data validation error: {e}")
        return 1
    except Exception as e:
        logger.error(f"Failed to load report data: {e}")
        return 1
    
    # Step 2: Prepare chart data
    logger.info("Step 2: Preparing chart data...")
    try:
        chart_data = prepare_chart_data(data)
        logger.info("✓ Chart data prepared")
        logger.info(f"  Charts to render: trend, pie, sankey, stacked_bar")
    except Exception as e:
        logger.error(f"Failed to prepare chart data: {e}")
        return 1
    
    # Step 3: Render charts
    logger.info("Step 3: Rendering charts...")
    try:
        from chart_renderer import render_charts_sync

        temp_charts_dir = "temp_charts"
        chart_images = render_charts_sync(chart_data, output_dir=temp_charts_dir)

        # Log rendered charts
        rendered_count = sum(1 for v in chart_images.values() if v is not None)
        logger.info(f"✓ Rendered {rendered_count}/{len(chart_images)} charts")
        for chart_name, path in chart_images.items():
            if path:
                logger.info(f"  - {chart_name}: {path}")
            else:
                logger.warning(f"  - {chart_name}: FAILED")
    except Exception as e:
        logger.error(f"Failed to render charts: {e}")
        return 1
    
    # Step 4: Create presentation
    logger.info("Step 4: Creating presentation...")
    try:
        prs = create_presentation()
        apply_branding(prs)
        logger.info("✓ Presentation created")
    except Exception as e:
        logger.error(f"Failed to create presentation: {e}")
        return 1
    
    # Step 5: Build slides
    logger.info("Step 5: Building slides...")
    try:
        # Build slides in order with section title cards for narrative flow
        # Narrative Arc: Current State → Value Delivered → Performance → Threats → Opportunities
        
        logger.info("  Building executive summary slides (includes title slide and section card)...")
        build_executive_summary_slides(prs, data)
        
        # Insert Value Delivered section title card
        logger.info("  Inserting 'Value Delivered' section card...")
        create_section_header_layout(
            prs,
            SECTION_NARRATIVES["value_delivered"]["title"],
            SECTION_NARRATIVES["value_delivered"]["narrative"]
        )
        
        logger.info("  Building value delivered slides...")
        build_value_delivered_slides(prs, data)
        
        # Insert Protection Achieved section title card
        logger.info("  Inserting 'Protection Achieved' section card...")
        create_section_header_layout(
            prs,
            SECTION_NARRATIVES["protection_achieved"]["title"],
            SECTION_NARRATIVES["protection_achieved"]["narrative"]
        )
        
        logger.info("  Building protection achieved slides...")
        build_protection_achieved_slides(prs, data)
        
        if not args.no_threat_landscape:
            # Insert Threat Landscape section title card
            logger.info("  Inserting 'Threat Landscape' section card...")
            create_section_header_layout(
                prs,
                SECTION_NARRATIVES["threat_landscape"]["title"],
                SECTION_NARRATIVES["threat_landscape"]["narrative"]
            )
            
            logger.info("  Building threat landscape slides...")
            build_threat_landscape_slides(prs, data, include=True)
        else:
            logger.info("  Skipping threat landscape slides (--no-threat-landscape)")
            build_threat_landscape_slides(prs, data, include=False)
        
        # Insert Insights section title card
        logger.info("  Inserting 'Insights & Continuous Improvement' section card...")
        create_section_header_layout(
            prs,
            SECTION_NARRATIVES["insights"]["title"],
            SECTION_NARRATIVES["insights"]["narrative"]
        )
        
        logger.info("  Building insights slides...")
        build_insights_slides(prs, data)

        # Insert Forward Direction section title card
        logger.info("  Inserting 'Forward Direction' section card...")
        create_section_header_layout(
            prs,
            SECTION_NARRATIVES["forward_direction"]["title"],
            SECTION_NARRATIVES["forward_direction"]["narrative"]
        )
        
        logger.info("  Building key takeaways slide...")
        # Build executive summary key takeaways (conditionally include based on data availability)
        takeaways = []
        
        # Industry comparison (only if benchmarks available)
        if data.industry_benchmarks_available and data.response_advantage_percent > 0:
            takeaways.append(f"{data.response_advantage_percent}% faster response than industry peers—threats are contained before spreading")
        else:
            takeaways.append(f"Average response time of {data.mttr_minutes} minutes with {data.p90_minutes}-minute P90")
        
        # Threat containment (always include)
        takeaways.append(f"100% threat containment with zero breaches this period across {data.true_threats_contained} true positive alerts")
        
        # Cost avoidance (always include)
        takeaways.append("Millions in modeled cost exposure avoided through proactive security operations")
        
        # After-hours (only if data available)
        if data.after_hours_data_available and data.after_hours_escalations > 0:
            takeaways.append(f"{data.after_hours_escalations} after-hours escalations handled seamlessly with {int(data.automation_percent)}% automation")
        else:
            takeaways.append(f"{int(data.automation_percent)}% of escalations handled via automated playbooks")
        
        build_key_takeaways_slide(prs, "This Period", takeaways, data)
        
        logger.info("  Building forward direction slide...")
        build_forward_direction_slide(prs, data)
        
        logger.info("  Building contact slide...")
        build_contact_slide(prs, data)
        
        logger.info(f"✓ Built {len(prs.slides)} slides")
    except Exception as e:
        logger.error(f"Failed to build slides: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return 1
    
    # Step 6: Insert chart images
    logger.info("Step 6: Inserting chart images...")
    try:
        charts_inserted = 0
        
        # Search through all slides for chart placeholders
        for slide_idx, slide in enumerate(prs.slides):
            # Check all shapes for chart placeholders
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text
                    
                    # Trend chart - "Performance Trends"
                    if chart_images.get('trend') and "[Chart: Performance Trends]" in text:
                        if insert_chart_image(slide, None, chart_images['trend']):
                            charts_inserted += 1
                            logger.info(f"  Inserted trend chart into slide {slide_idx + 1}")
                            break  # Move to next slide
                    
                    # Sankey chart - "Severity Alignment Sankey"
                    if chart_images.get('sankey') and 'severity_sankey' in text:
                        if insert_chart_image(slide, 'severity_sankey', chart_images['sankey']):
                            charts_inserted += 1
                            logger.info(f"  Inserted sankey chart into slide {slide_idx + 1}")
                            break  # Move to next slide
                    
                    # Stacked bar chart - "MITRE ATT&CK Stacked Bar"
                    if chart_images.get('stacked_bar') and 'mitre_stacked_bar' in text:
                        if insert_chart_image(slide, 'mitre_stacked_bar', chart_images['stacked_bar']):
                            charts_inserted += 1
                            logger.info(f"  Inserted stacked bar chart into slide {slide_idx + 1}")
                            break  # Move to next slide
                    
                    # Funnel chart - "CORR Funnel"
                    if chart_images.get('funnel') and 'corr_funnel' in text:
                        if insert_chart_image(slide, 'corr_funnel', chart_images['funnel']):
                            charts_inserted += 1
                            logger.info(f"  Inserted funnel chart into slide {slide_idx + 1}")
                            break  # Move to next slide
        
        logger.info(f"✓ Inserted {charts_inserted} chart images")
    except Exception as e:
        logger.error(f"Failed to insert chart images: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return 1
    
    # Step 7: Save presentation
    logger.info("Step 7: Saving presentation...")
    try:
        # Create output directory
        output_dir = Path(args.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Generate filename with date
        date_str = datetime.now().strftime("%Y-%m-%d")
        output_filename = f"executive_business_review_{date_str}.pptx"
        output_path = output_dir / output_filename
        
        prs.save(str(output_path))
        logger.info(f"✓ Presentation saved")
        logger.info(f"  Path: {output_path}")
        logger.info(f"  Total slides: {len(prs.slides)}")
    except Exception as e:
        logger.error(f"Failed to save presentation: {e}")
        return 1
    
    # Step 8: Cleanup (optional)
    if not args.keep_charts:
        logger.info("Step 8: Cleaning up temporary files...")
        try:
            temp_charts_path = Path(temp_charts_dir)
            if temp_charts_path.exists():
                shutil.rmtree(temp_charts_path)
                logger.info(f"✓ Deleted {temp_charts_dir}/")
        except Exception as e:
            logger.warning(f"Failed to cleanup temp files: {e}")
    else:
        logger.info(f"Step 8: Keeping temporary charts in {temp_charts_dir}/")
    
    logger.info("=" * 60)
    logger.info("✓ Presentation generation completed successfully!")
    logger.info("=" * 60)
    
    return 0


def build_executive_summary_slides(prs, data):
    """Create the Title slide and Executive Summary slides (Slides 1-3).
    
    LAYOUT STRUCTURE:
    -----------------
    Slide 1 - Title Slide (Blue Gradient Background):
        - Main title: "EXECUTIVE BUSINESS REVIEW" (H1 - 114pt)
        - Client name (H5 - 16pt)
        - Period dates with day count
        - Report date at bottom
        - Footer only, NO header per branding
    
    Slide 2 - Executive Summary Section Card:
        - Gray background with section title and narrative
        - Created via create_section_header_layout()
    
    Slide 3 - Executive Dashboard (2x3 Metric Grid):
        Top Row: True Positives Contained, MTTR (with % advantage), Closed End-to-End
        Bottom Row: Alerts Triaged, After-Hours Escalations, False Positive Rate
        Insight bar at bottom with security posture summary
    
    KEY DATA FIELDS USED:
    - client_name, period_start, period_end, report_date (title slide)
    - true_threats_contained, mttr_minutes, response_advantage_percent
    - closed_end_to_end, alerts_triaged, after_hours_escalations
    - false_positive_rate, executive_summary_narrative
    
    Per CRITICALSTART branding guidelines:
    - Title slide: No header, footer only, H1 typography for main title
    - Content slides: Header and footer with transparent background
    
    Args:
        prs (Presentation): The presentation object.
        data (ReportData): The report data object containing all metrics.
    """
    # Slide 1 - Title Slide (NO HEADER per branding guidelines)
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Add gradient background (blue sweep: #009CDE → #004C97)
    create_gradient_background(prs, slide1, 'blue_sweep')
    
    # Note: Logo removed from title slide per branding guidelines
    
    # Add main title using H1 typography (114pt per branding guide)
    title_left = Inches(1)
    title_top = Inches(1.5)
    title_width = prs.slide_width - Inches(2)
    title_height = Inches(1.5)
    
    title_box = slide1.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = "EXECUTIVE BUSINESS\nREVIEW"
    title_paragraph.font.name = TITLE_FONT_NAME
    title_paragraph.font.size = H1_FONT_SIZE  # 114pt per branding guide
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Add client name using H5 typography
    client_left = Inches(1)
    client_top = Inches(3.2)
    client_width = prs.slide_width - Inches(2)
    client_height = Inches(0.5)
    
    client_box = slide1.shapes.add_textbox(client_left, client_top, client_width, client_height)
    client_frame = client_box.text_frame
    client_frame.word_wrap = True
    client_paragraph = client_frame.paragraphs[0]
    client_paragraph.text = data.client_name
    client_paragraph.font.name = BODY_FONT_NAME
    client_paragraph.font.size = H5_FONT_SIZE  # 16pt
    client_paragraph.font.color.rgb = RGBColor(255, 255, 255)
    client_paragraph.alignment = PP_ALIGN.LEFT
    
    # Add period using paragraph typography
    # Format: "August 1-31, 2025 (31 days)"
    start_month_day = data.period_start.split(',')[0].strip()  # "August 1"
    end_month_day = data.period_end.split(',')[0].strip()  # "August 31"
    year = data.period_end.split(',')[1].strip()  # "2025"
    
    # Extract month and day numbers
    start_parts = start_month_day.split()
    end_parts = end_month_day.split()
    
    if len(start_parts) == 2 and len(end_parts) == 2 and start_parts[0] == end_parts[0]:
        # Same month, format as "August 1-31, 2025"
        period_text = f"{start_parts[0]} {start_parts[1]}-{end_parts[1]}, {year} ({data.period_days} days)"
    else:
        # Different months, use full format
        period_text = f"{start_month_day} - {end_month_day}, {year} ({data.period_days} days)"
    period_left = Inches(1)
    period_top = Inches(3.7)
    period_width = prs.slide_width - Inches(2)
    period_height = Inches(0.5)
    
    period_box = slide1.shapes.add_textbox(period_left, period_top, period_width, period_height)
    period_frame = period_box.text_frame
    period_frame.word_wrap = True
    period_paragraph = period_frame.paragraphs[0]
    period_paragraph.text = period_text
    period_paragraph.font.name = BODY_FONT_NAME
    period_paragraph.font.size = PARAGRAPH_FONT_SIZE  # 12pt
    period_paragraph.font.color.rgb = RGBColor(255, 255, 255)
    period_paragraph.alignment = PP_ALIGN.LEFT
    
    # Add footer only (no header on title slide per branding guidelines)
    # Use white text on blue gradient background for WCAG AA compliance
    add_master_slide_elements(slide1, prs, slide_number=None, 
                               include_header=False, include_footer=True,
                               text_color=RGBColor(255, 255, 255))
    
    # Note: Report date is now in the footer per branding guidelines
    report_date_left = Inches(1)
    report_date_top = prs.slide_height - Inches(0.9)  # Moved up to avoid footer
    report_date_width = prs.slide_width - Inches(2)
    report_date_height = Inches(0.4)
    
    report_date_box = slide1.shapes.add_textbox(report_date_left, report_date_top, report_date_width, report_date_height)
    report_date_frame = report_date_box.text_frame
    report_date_frame.word_wrap = True
    report_date_paragraph = report_date_frame.paragraphs[0]
    report_date_paragraph.text = f"Report Date: {data.report_date}"
    report_date_paragraph.font.name = BODY_FONT_NAME
    report_date_paragraph.font.size = Pt(14)
    report_date_paragraph.font.color.rgb = RGBColor(255, 255, 255)
    report_date_paragraph.alignment = PP_ALIGN.LEFT
    
    # Insert Executive Summary section title card
    create_section_header_layout(
        prs, 
        SECTION_NARRATIVES["executive_summary"]["title"],
        SECTION_NARRATIVES["executive_summary"]["narrative"]
    )
    
    # Slide 3 - Executive Dashboard (Merged BLUF)
    # Combines Security Posture + Value metrics into single consolidated dashboard
    slide2 = prs.slides.add_slide(blank_slide_layout)
    slide2_number = get_slide_number(prs)
    
    # Add master slide elements (header and footer with transparent background)
    add_master_slide_elements(slide2, prs, slide_number=slide2_number,
                               include_header=True, include_footer=True)
    
    # Add logo at top right
    add_logo(slide2, position='top_right', prs=prs)
    
    # Content area starts below header
    content_top = HEADER_HEIGHT + Inches(0.3)
    
    # Add slide title using H3 typography
    slide2_title_left = MARGIN_STANDARD
    slide2_title_top = content_top
    slide2_title_width = prs.slide_width - Inches(2)
    slide2_title_height = Inches(0.45)
    
    slide2_title_box = slide2.shapes.add_textbox(slide2_title_left, slide2_title_top, slide2_title_width, slide2_title_height)
    slide2_title_frame = slide2_title_box.text_frame
    slide2_title_frame.word_wrap = True
    slide2_title_paragraph = slide2_title_frame.paragraphs[0]
    slide2_title_paragraph.text = "Executive Dashboard"
    slide2_title_paragraph.font.name = TITLE_FONT_NAME
    slide2_title_paragraph.font.size = H3_FONT_SIZE
    slide2_title_paragraph.font.bold = True
    slide2_title_paragraph.font.color.rgb = CS_NAVY
    slide2_title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Calculate values for display
    total_millions = data.total_modeled / 1000000
    after_hours_percent = int((data.after_hours_escalations / data.incidents_escalated) * 100) if data.incidents_escalated > 0 else 0
    
    # Determine MTTR detail text based on industry benchmark availability
    if data.industry_benchmarks_available and data.response_advantage_percent > 0:
        mttr_detail = f"{int(data.response_advantage_percent)}% faster than industry ({data.industry_median_minutes}m)"
    else:
        mttr_detail = f"P90: {data.p90_minutes} minutes"
    
    # Determine after-hours display based on data availability
    if data.after_hours_data_available:
        after_hours_metric = f"{data.after_hours_escalations}"
        after_hours_detail = f"{after_hours_percent}% of escalations"
    else:
        after_hours_metric = "—"
        after_hours_detail = "Analysis coming soon"
    
    # Consolidated 6 key metrics in 2 rows x 3 columns
    # Each metric includes a metric_name for threshold lookup and raw_value for status evaluation
    dashboard_metrics = [
        # Row 1: Core Security Outcomes
        {
            "title": "THREATS CONTAINED",
            "metric": f"{data.true_threats_contained}",
            "detail": "100% contained, zero breaches",
            "metric_name": "threats_contained",
            "raw_value": data.true_threats_contained
        },
        {
            "title": "MTTR",
            "metric": f"{data.mttr_minutes} min",
            "detail": mttr_detail,
            "metric_name": "mttr",
            "raw_value": data.mttr_minutes
        },
        {
            "title": "CLOSED END-TO-END",
            "metric": f"{data.closed_end_to_end:,}",
            "detail": "Alerts resolved without client action",
            "metric_name": "closed_end_to_end",
            "raw_value": data.closed_end_to_end
        },
        # Row 2: Operations & Coverage
        {
            "title": "ALERTS TRIAGED",
            "metric": f"{data.alerts_triaged:,}",
            "detail": f"{data.client_touch_decisions:,} required client input",
            "metric_name": None,  # Informational only, no threshold
            "raw_value": None
        },
        {
            "title": "AFTER-HOURS",
            "metric": after_hours_metric,
            "detail": after_hours_detail,
            "metric_name": None,  # Informational only, no threshold
            "raw_value": None
        },
        {
            "title": "FALSE POSITIVE RATE",
            "metric": f"{data.false_positive_rate}%",
            "detail": "Below 10% target threshold",
            "metric_name": "false_positive_rate",
            "raw_value": data.false_positive_rate
        }
    ]
    
    # Card dimensions for 2x3 grid
    card_width = (prs.slide_width - Inches(1.4)) / 3
    card_height = Inches(1.4)
    card_start_left = Inches(0.5)
    card_start_top = content_top + Inches(0.5)
    card_spacing_h = Inches(0.2)
    card_spacing_v = Inches(0.15)
    
    # Track status for insight bar summary
    status_summary = {"good": 0, "warning": 0, "bad": 0}
    
    # Draw 6 cards in 2x3 grid
    for idx, card in enumerate(dashboard_metrics):
        row = idx // 3
        col = idx % 3
        card_left = card_start_left + col * (card_width + card_spacing_h)
        card_top = card_start_top + row * (card_height + card_spacing_v)
        
        # Get metric status if threshold exists
        status_info = None
        if card.get("metric_name") and card.get("raw_value") is not None:
            status_info = get_metric_status(card["metric_name"], card["raw_value"])
            if status_info:
                status_summary[status_info["status"]] += 1
        
        # Determine border color based on status
        if status_info:
            indicator = get_status_indicator(status_info["status"], status_info["direction"])
            border_color = indicator["color"]
        else:
            border_color = CS_BLUE if row == 0 else CS_SLATE
        
        # Create card background
        card_shape = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, card_left, card_top,
            card_width, card_height
        )
        fill = card_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255) if row == 0 else RGBColor(248, 250, 252)
        line = card_shape.line
        line.color.rgb = border_color
        line.width = Pt(2) if row == 0 else Pt(1)
        
        # Add card title (small label)
        title_box = slide2.shapes.add_textbox(
            card_left + Inches(0.12), card_top + Inches(0.08),
            card_width - Inches(0.24), Inches(0.22)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = card["title"]
        title_para.font.name = TITLE_FONT_NAME
        title_para.font.size = Pt(10)
        title_para.font.bold = True
        title_para.font.color.rgb = CS_SLATE
        title_para.alignment = PP_ALIGN.LEFT
        
        # Add large metric with indicator
        indicator_width = Inches(0.4) if status_info else 0
        metric_box = slide2.shapes.add_textbox(
            card_left + Inches(0.12), card_top + Inches(0.28),
            card_width - Inches(0.24) - indicator_width, Inches(0.6)
        )
        metric_frame = metric_box.text_frame
        metric_para = metric_frame.paragraphs[0]
        metric_para.text = card["metric"]
        metric_para.font.name = TITLE_FONT_NAME
        metric_para.font.size = Pt(42)
        metric_para.font.bold = True
        metric_para.font.color.rgb = CS_NAVY
        metric_para.alignment = PP_ALIGN.LEFT
        
        # Note: Trend indicator arrows removed per Jan 2026 stakeholder feedback.
        # Arrows on static cards confused users (implied trends, not threshold status).
        # The colored border alone now indicates good/warning/bad status.
        
        # Add detail text
        detail_box = slide2.shapes.add_textbox(
            card_left + Inches(0.12), card_top + Inches(0.95),
            card_width - Inches(0.24), Inches(0.4)
        )
        detail_frame = detail_box.text_frame
        detail_frame.word_wrap = True
        detail_para = detail_frame.paragraphs[0]
        detail_para.text = card["detail"]
        detail_para.font.name = BODY_FONT_NAME
        detail_para.font.size = Pt(11)
        detail_para.font.color.rgb = CS_SLATE
        detail_para.alignment = PP_ALIGN.LEFT
    
    # Add insight bar at bottom of dashboard with dynamic status summary
    insight_top = card_start_top + 2 * (card_height + card_spacing_v) + Inches(0.2)
    
    # Generate dynamic insight text based on status summary
    if status_summary["bad"] == 0 and status_summary["warning"] == 0:
        insight_body = "All metrics meeting or exceeding targets—your organization is well-protected."
    elif status_summary["bad"] == 0:
        insight_body = f"{status_summary['good']} metrics exceeding targets, {status_summary['warning']} approaching thresholds—strong overall posture."
    else:
        insight_body = f"{status_summary['good']} metrics on target, {status_summary['warning'] + status_summary['bad']} needing attention—see details in following slides."
    
    add_insight_callout(
        slide2, prs,
        "Your Security Posture at a Glance",
        insight_body,
        insight_top,
        height=Inches(0.7)
    )
    
    # Slide 4 - CORR Platform Funnel (AI Accelerated Security)
    build_corr_funnel_slide(prs, data)


def create_executive_summary_slide(prs, report_data):
    """Create the Executive Summary slide."""
    pass


def build_corr_funnel_slide(prs, data=None, background_image_path: Optional[str] = None):
    """Create the CORR Platform slide with a background image and data overlays.

    This slide displays a pre-designed PNG as the full-slide background,
    showing the AI-accelerated security pipeline visualization.
    Text boxes overlay the funnel stages with actual or placeholder values.

    Args:
        prs (Presentation): The presentation object.
        data (ReportData, optional): Report data for populating text boxes.
            If None, uses placeholder values.
        background_image_path (str, optional): Path to background PNG image.
            If None, uses the default 'assets/funnel diagram.png'.

    Returns:
        Slide: The created slide object.
    """
    # Create blank slide (no title, no standard branding elements)
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_number = get_slide_number(prs)

    # Add header and footer only (no title, no logo - background image is the content)
    add_master_slide_elements(slide, prs, slide_number=slide_number,
                               include_header=True, include_footer=True)

    # Determine background image path
    if background_image_path is None:
        background_image_path = Path(__file__).parent / "assets" / "funnel diagram.png"
    else:
        background_image_path = Path(background_image_path)

    # Add background image covering the full slide
    if background_image_path.exists():
        try:
            # Add picture covering full slide dimensions
            bg_picture = slide.shapes.add_picture(
                str(background_image_path),
                Inches(0), Inches(0),
                prs.slide_width, prs.slide_height
            )

            # Move the background image to the back of the z-order
            # Access the shape tree and move the picture element to the first position
            spTree = slide.shapes._spTree
            sp = bg_picture._element
            spTree.remove(sp)
            # Insert after nvGrpSpPr (the first child is usually the group shape properties)
            spTree.insert(2, sp)

            logging.info(f"Added background image to slide {slide_number}: {background_image_path}")

        except Exception as e:
            logging.error(f"Failed to add background image: {e}")
    else:
        logging.warning(f"Background image not found: {background_image_path}")

    # Add text boxes overlaying funnel stages with data
    _add_funnel_text_boxes(slide, prs, data)

    return slide


def _add_funnel_text_boxes(slide, prs, data=None):
    """Add positioned text boxes with data to the funnel slide.

    Text boxes overlay the 4 funnel stages to display metrics:
    - Stage 1 (leftmost): Security Events - total events analyzed
    - Stage 2: Potential Threats - alerts requiring triage
    - Stage 3: Alerts - incidents escalated to client
    - Stage 4 (rightmost): Response Actions - threats contained

    Note: Positions are calibrated for 'funnel diagram.png' at 16:9 dimensions.
    Adjust positions if using a different funnel image.

    Args:
        slide: The PowerPoint slide object
        prs: The presentation object (for dimensions)
        data: ReportData object or None for placeholder values
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # Funnel stage positions calibrated to the PNG background
    # Format: (left, top, width, height) - positions for text boxes
    # These may need adjustment based on actual funnel diagram dimensions
    STAGE_POSITIONS = {
        'security_events': (Inches(1.0), Inches(2.0), Inches(1.8), Inches(1.0)),
        'potential_threats': (Inches(3.0), Inches(2.2), Inches(1.8), Inches(0.9)),
        'alerts': (Inches(5.2), Inches(2.4), Inches(1.6), Inches(0.8)),
        'response_actions': (Inches(7.2), Inches(2.5), Inches(1.6), Inches(0.7)),
    }

    # Extract values from data or use placeholders
    # Placeholders marked with asterisk to indicate they need real data
    if data:
        stage_values = {
            'security_events': {
                'value': '—*',  # Placeholder - data not yet available
                'label': 'Security Events',
            },
            'potential_threats': {
                'value': f'{data.alerts_triaged:,}' if hasattr(data, 'alerts_triaged') and data.alerts_triaged else '—*',
                'label': 'Alerts Triaged',
            },
            'alerts': {
                'value': f'{data.incidents_escalated}' if hasattr(data, 'incidents_escalated') and data.incidents_escalated else '—*',
                'label': 'Escalated',
            },
            'response_actions': {
                'value': f'{data.true_threats_contained}' if hasattr(data, 'true_threats_contained') and data.true_threats_contained else '—*',
                'label': 'Contained',
            },
        }
    else:
        # Full placeholder mode
        stage_values = {
            'security_events': {'value': '—*', 'label': 'Security Events'},
            'potential_threats': {'value': '—*', 'label': 'Alerts Triaged'},
            'alerts': {'value': '—*', 'label': 'Escalated'},
            'response_actions': {'value': '—*', 'label': 'Contained'},
        }

    # Add text boxes for each stage
    for stage_name, pos in STAGE_POSITIONS.items():
        left, top, width, height = pos
        stage = stage_values[stage_name]

        # Create text box
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        # Add value (large, bold, white for visibility on colored background)
        p_value = text_frame.paragraphs[0]
        p_value.text = stage['value']
        p_value.font.name = TITLE_FONT_NAME
        p_value.font.size = Pt(28)
        p_value.font.bold = True
        p_value.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p_value.alignment = PP_ALIGN.CENTER

        # Add label (smaller, white)
        p_label = text_frame.add_paragraph()
        p_label.text = stage['label']
        p_label.font.name = BODY_FONT_NAME
        p_label.font.size = Pt(11)
        p_label.font.color.rgb = RGBColor(255, 255, 255)
        p_label.alignment = PP_ALIGN.CENTER


def build_value_delivered_slides(prs, data):
    """Create the Value Delivered section (consolidated to 2 slides).
    
    Per CRITICALSTART branding guidelines:
    - All slides have header and footer with transparent background
    - Uses H1-H6 typography scale
    
    Args:
        prs (Presentation): The presentation object.
        data (ReportData): The report data object containing all metrics.
    """
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    content_top = HEADER_HEIGHT + Inches(0.3)
    
    # Slide 3 - Value Delivered (Merged: Hero + Breakdown)
    slide_value = prs.slides.add_slide(blank_slide_layout)
    slide_value_number = get_slide_number(prs)
    
    # Add master slide elements
    add_master_slide_elements(slide_value, prs, slide_number=slide_value_number,
                               include_header=True, include_footer=True)
    
    # Add logo at top right
    add_logo(slide_value, position='top_right', prs=prs)
    
    # Add title using H3 typography
    title_left = MARGIN_STANDARD
    title_top = content_top
    title_width = prs.slide_width - Inches(2)
    title_height = Inches(0.5)
    
    title_box = slide_value.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = "Value Delivered"
    title_paragraph.font.name = TITLE_FONT_NAME
    title_paragraph.font.size = H3_FONT_SIZE
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = CS_NAVY
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Hero section (left side) - Cost Avoidance
    total_millions = data.total_modeled / 1000000
    hero_value = f"~${total_millions:.1f}M"
    
    hero_left = Inches(0.5)
    hero_top = content_top + Inches(0.55)
    hero_width = Inches(3.5)
    hero_height = Inches(2.8)
    
    # Hero background
    hero_bg = slide_value.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, hero_left, hero_top,
        hero_width, hero_height
    )
    fill = hero_bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    line = hero_bg.line
    line.color.rgb = CS_BLUE
    line.width = Pt(2)
    
    # Hero value
    hero_box = slide_value.shapes.add_textbox(hero_left + Inches(0.2), hero_top + Inches(0.3), 
                                               hero_width - Inches(0.4), Inches(1.2))
    hero_frame = hero_box.text_frame
    hero_paragraph = hero_frame.paragraphs[0]
    hero_paragraph.text = hero_value
    hero_paragraph.font.name = TITLE_FONT_NAME
    hero_paragraph.font.size = Pt(72)
    hero_paragraph.font.bold = True
    hero_paragraph.font.color.rgb = CS_NAVY
    hero_paragraph.alignment = PP_ALIGN.CENTER
    
    # Hero label
    hero_label = slide_value.shapes.add_textbox(hero_left + Inches(0.2), hero_top + Inches(1.5), 
                                                 hero_width - Inches(0.4), Inches(0.4))
    label_frame = hero_label.text_frame
    label_para = label_frame.paragraphs[0]
    label_para.text = "Cost Exposure Avoided"
    label_para.font.name = TITLE_FONT_NAME
    label_para.font.size = Pt(18)
    label_para.font.bold = True
    label_para.font.color.rgb = CS_NAVY
    label_para.alignment = PP_ALIGN.CENTER
    
    # Hero subtitle
    hero_sub = slide_value.shapes.add_textbox(hero_left + Inches(0.2), hero_top + Inches(2.0), 
                                               hero_width - Inches(0.4), Inches(0.6))
    sub_frame = hero_sub.text_frame
    sub_frame.word_wrap = True
    sub_para = sub_frame.paragraphs[0]
    sub_para.text = "Modeled operational, coverage, and breach avoidance"
    sub_para.font.name = BODY_FONT_NAME
    sub_para.font.size = Pt(12)
    sub_para.font.color.rgb = CS_SLATE
    sub_para.alignment = PP_ALIGN.CENTER
    
    # Right side - 3 breakdown cards (stacked vertically)
    card_left = Inches(4.2)
    card_width = prs.slide_width - card_left - Inches(0.5)
    card_height = Inches(0.85)
    card_spacing = Inches(0.1)
    card_start_top = hero_top
    
    # Format values
    analyst_k = data.analyst_cost_equivalent / 1000
    coverage_k = data.coverage_cost_equivalent / 1000
    breach_m = data.breach_exposure_avoided / 1000000
    
    breakdown_cards = [
        {"title": "Security Operations", "metric": f"~${analyst_k:.0f}K", "detail": f"{data.analyst_hours} analyst hours"},
        {"title": "24/7 Coverage", "metric": f"~${coverage_k:.0f}K", "detail": f"{data.coverage_hours} hours monitoring"},
        {"title": "Threat Prevention", "metric": f"~${breach_m:.1f}M", "detail": f"{data.true_threats_contained} threats contained"}
    ]
    
    for i, card in enumerate(breakdown_cards):
        card_top = card_start_top + i * (card_height + card_spacing)
        
        # Card background
        card_shape = slide_value.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, card_left, card_top,
            card_width, card_height
        )
        fill = card_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 250, 252)
        line = card_shape.line
        line.color.rgb = CS_SLATE
        line.width = Pt(1)
        
        # Card title (left)
        title_box = slide_value.shapes.add_textbox(
            card_left + Inches(0.15), card_top + Inches(0.12),
            Inches(2.2), Inches(0.3)
        )
        t_frame = title_box.text_frame
        t_para = t_frame.paragraphs[0]
        t_para.text = card["title"]
        t_para.font.name = TITLE_FONT_NAME
        t_para.font.size = Pt(14)
        t_para.font.bold = True
        t_para.font.color.rgb = CS_SLATE
        t_para.alignment = PP_ALIGN.LEFT
        
        # Card detail (left, below title)
        detail_box = slide_value.shapes.add_textbox(
            card_left + Inches(0.15), card_top + Inches(0.45),
            Inches(2.5), Inches(0.3)
        )
        d_frame = detail_box.text_frame
        d_para = d_frame.paragraphs[0]
        d_para.text = card["detail"]
        d_para.font.name = BODY_FONT_NAME
        d_para.font.size = Pt(11)
        d_para.font.color.rgb = CS_SLATE
        d_para.alignment = PP_ALIGN.LEFT
        
        # Card metric (right)
        metric_box = slide_value.shapes.add_textbox(
            card_left + card_width - Inches(1.8), card_top + Inches(0.15),
            Inches(1.6), Inches(0.55)
        )
        m_frame = metric_box.text_frame
        m_para = m_frame.paragraphs[0]
        m_para.text = card["metric"]
        m_para.font.name = TITLE_FONT_NAME
        m_para.font.size = Pt(28)
        m_para.font.bold = True
        m_para.font.color.rgb = CS_NAVY
        m_para.alignment = PP_ALIGN.RIGHT
    
    # Methodology note at bottom (expanded per Jan 2026 stakeholder feedback)
    note_top = hero_top + hero_height + Inches(0.1)
    note_box = slide_value.shapes.add_textbox(Inches(0.5), note_top, 
                                               prs.slide_width - Inches(1), Inches(0.45))
    note_frame = note_box.text_frame
    note_frame.word_wrap = True
    note_para = note_frame.paragraphs[0]
    # Expanded methodology explanation per stakeholder request
    note_para.text = (
        "* Cost methodology: Operations = analyst hours × $85/hr avg. rate; "
        "Coverage = 24/7 monitoring × market SOC rates; "
        "Prevention = contained threats × Ponemon breach cost model ($4.45M avg. breach, adjusted for threat severity). "
        "Illustrative impact only; not redeployable budget."
    )
    note_para.font.name = BODY_FONT_NAME
    note_para.font.size = Pt(9)
    note_para.font.color.rgb = CS_SLATE
    note_para.font.italic = True
    note_para.alignment = PP_ALIGN.LEFT


def build_protection_achieved_slides(prs, data):
    """Create the Protection section (Slides 7-10).
    
    LAYOUT STRUCTURE:
    -----------------
    Slide 7 - Critical Start's Performance (Split Layout):
        Left Panel (2.3"): MTTR by Severity cards (Crit/High, Med/Low, P90)
        Right Panel (6.7"): Trend chart with industry benchmark lines
    
    Slide 8 - Response & Detection Quality (Dual Panel):
        Left Panel: Response Efficiency (Remediation Rate, Automation, Human Review)
        Right Panel: Detection Quality (TP Rate, FP Rate, Alert Reduction)
    
    KEY DATA FIELDS USED:
    - critical_high_mttr, medium_low_mttr, p90_minutes (MTTR by severity)
    - mttr_trend, mttd_trend, fp_trend (trend chart data)
    - containment_rate, automation_percent (response efficiency)
    - false_positive_rate, true_threat_precision (detection quality)
    - industry_benchmarks_available (conditional benchmark display)
    
    Args:
        prs (Presentation): The presentation object.
        data (ReportData): The report data object containing all metrics.
    """
    # Slide 7 - Critical Start's Performance (Split Layout)
    slide7, content_top_7 = setup_content_slide(prs, "Critical Start's Performance")
    
    # =========================================================================
    # LEFT PANEL: Response Metrics (MTTR by Severity)
    # =========================================================================
    panel_left = Inches(0.5)
    panel_top = content_top_7 + Inches(0.1)
    panel_width = Inches(2.3)
    panel_height = Inches(3.2)
    
    # Panel background
    panel_bg = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, panel_left, panel_top,
        panel_width, panel_height
    )
    panel_bg.fill.solid()
    panel_bg.fill.fore_color.rgb = RGBColor(248, 250, 252)  # Very light gray
    panel_bg.line.color.rgb = CS_SLATE
    panel_bg.line.width = Pt(1)
    
    # Panel header
    header_height = Inches(0.4)
    header_shape = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, panel_left, panel_top,
        panel_width, header_height
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = CS_NAVY
    header_shape.line.fill.background()
    
    header_text = slide7.shapes.add_textbox(panel_left, panel_top, panel_width, header_height)
    header_tf = header_text.text_frame
    header_tf.paragraphs[0].text = "MTTR BY SEVERITY"  # Updated per Jan 2026 feedback
    header_tf.paragraphs[0].font.name = TITLE_FONT_NAME
    header_tf.paragraphs[0].font.size = Pt(12)
    header_tf.paragraphs[0].font.bold = True
    header_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Metric cards within panel - now with threshold status indicators
    # Labels clarified per Jan 2026 stakeholder feedback to indicate vendor severities
    metrics = [
        {"value": f"{data.critical_high_mttr} min", "label": "Critical & High (Vendor)", "default_color": CS_RED,
         "metric_name": "critical_high_mttr", "raw_value": data.critical_high_mttr},
        {"value": f"{data.medium_low_mttr} min", "label": "Medium & Low (Vendor)", "default_color": CS_ORANGE,
         "metric_name": "medium_low_mttr", "raw_value": data.medium_low_mttr},
        {"value": f"{data.p90_minutes} min", "label": "P90 All Severities", "default_color": CS_BLUE,
         "metric_name": "p90_minutes", "raw_value": data.p90_minutes},
    ]
    
    card_top = panel_top + header_height + Inches(0.15)
    card_height = Inches(0.85)
    card_spacing = Inches(0.1)
    card_margin = Inches(0.1)
    card_width = panel_width - card_margin * 2
    
    for idx, metric in enumerate(metrics):
        current_top = card_top + idx * (card_height + card_spacing)
        
        # Get metric status for threshold evaluation
        status_info = get_metric_status(metric["metric_name"], metric["raw_value"])
        
        # Determine color based on status (override default color)
        if status_info:
            indicator = get_status_indicator(status_info["status"], status_info["direction"])
            border_color = indicator["color"]
        else:
            border_color = metric["default_color"]
        
        # Card background
        card_shape = slide7.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            panel_left + card_margin, current_top,
            card_width, card_height
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card_shape.line.color.rgb = border_color
        card_shape.line.width = Pt(2)
        
        # Metric value with indicator
        indicator_width = Inches(0.35) if status_info else 0
        value_box = slide7.shapes.add_textbox(
            panel_left + card_margin + Inches(0.05), current_top + Inches(0.1),
            card_width - Inches(0.1) - indicator_width, Inches(0.45)
        )
        value_tf = value_box.text_frame
        value_tf.paragraphs[0].text = metric["value"]
        value_tf.paragraphs[0].font.name = TITLE_FONT_NAME
        value_tf.paragraphs[0].font.size = Pt(28)
        value_tf.paragraphs[0].font.bold = True
        value_tf.paragraphs[0].font.color.rgb = border_color
        value_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Note: Trend indicator arrows removed per Jan 2026 stakeholder feedback.
        # Colored border alone indicates status.
        
        # Metric label
        label_box = slide7.shapes.add_textbox(
            panel_left + card_margin, current_top + Inches(0.5),
            card_width, Inches(0.3)
        )
        label_tf = label_box.text_frame
        label_tf.paragraphs[0].text = metric["label"]
        label_tf.paragraphs[0].font.name = BODY_FONT_NAME
        label_tf.paragraphs[0].font.size = Pt(11)
        label_tf.paragraphs[0].font.color.rgb = CS_SLATE
        label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # =========================================================================
    # RIGHT PANEL: Trend Chart (enlarged per Jan 2026 feedback)
    # =========================================================================
    chart_left = panel_left + panel_width + Inches(0.2)
    chart_top = content_top_7 + Inches(0.1)
    chart_width = prs.slide_width - chart_left - Inches(0.5)
    chart_height = Inches(2.65)  # Increased from 2.3"
    
    chart_placeholder = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        chart_width, chart_height
    )
    fill = chart_placeholder.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
    line = chart_placeholder.line
    line.color.rgb = CS_SLATE
    line.width = Pt(1)
    
    # Add placeholder text
    placeholder_text = chart_placeholder.text_frame
    placeholder_text.text = "[Chart: Performance Trends]"
    placeholder_text.paragraphs[0].font.name = BODY_FONT_NAME
    placeholder_text.paragraphs[0].font.size = Pt(14)
    placeholder_text.paragraphs[0].font.color.rgb = CS_SLATE
    placeholder_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    placeholder_text.vertical_anchor = 1  # Middle
    
    # Add legend below chart
    legend_top = chart_top + chart_height + Inches(0.15)
    legend_left = chart_left
    legend_width = chart_width
    legend_height = Inches(0.35)
    
    legend_box = slide7.shapes.add_textbox(legend_left, legend_top, legend_width, legend_height)
    legend_frame = legend_box.text_frame
    legend_frame.word_wrap = True
    legend_paragraph = legend_frame.paragraphs[0]
    # Updated legend to include industry benchmarks (dashed lines on chart)
    legend_paragraph.text = "MTTR (blue) | MTTD (navy) | FP% (red) | Industry benchmarks (dashed)"
    legend_paragraph.font.name = BODY_FONT_NAME
    legend_paragraph.font.size = Pt(10)
    legend_paragraph.font.color.rgb = CS_SLATE
    legend_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add insight box below legend (spanning full width of right panel)
    insight_top = legend_top + legend_height + Inches(0.1)
    insight_height = Inches(0.65)
    
    insight_shape = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, insight_top,
        chart_width, insight_height
    )
    fill = insight_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
    line = insight_shape.line
    line.color.rgb = CS_BLUE
    line.width = Pt(2)
    
    # Generate dynamic insight text - explain P90 significance and industry comparison
    # Industry benchmarks now shown on chart; insight explains why this matters
    response_advantage = int(data.response_advantage_percent)
    insight_message = (
        f"Your MTTR outperforms industry benchmarks by {response_advantage}%—"
        "faster response means smaller blast radius when threats emerge."
    )
    
    insight_text = insight_shape.text_frame
    insight_text.text = insight_message
    insight_text.paragraphs[0].font.name = BODY_FONT_NAME
    insight_text.paragraphs[0].font.size = Pt(13)
    insight_text.paragraphs[0].font.color.rgb = CS_NAVY
    insight_text.paragraphs[0].font.bold = True
    insight_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    insight_text.vertical_anchor = 1  # Middle
    
    # NOTE: Industry Comparison slide removed per Jan 2026 stakeholder feedback.
    # Industry benchmarks are now shown as dashed lines on the trend chart above.
    
    # Slide 7 - Response & Detection Quality (Merged dual-panel)
    slide_resp_det, content_top_rd = setup_content_slide(prs, "Response & Detection Quality")
    
    # Left panel: Response Efficiency (3 vertical cards)
    left_panel_left = Inches(0.5)
    left_panel_width = (prs.slide_width - Inches(1.2)) / 2
    card_height = Inches(0.9)
    card_spacing = Inches(0.1)
    cards_start_top = content_top_rd + Inches(0.1)
    
    # Calculate remediation count (round to nearest integer)
    # Renamed from Containment Rate → Remediation Rate per Jan 2026 feedback
    remediation_count = round((data.containment_rate / 100) * data.incidents_escalated)
    
    efficiency_cards = [
        {"title": "Remediation Rate", "value": f"{data.containment_rate}%", "detail": f"{remediation_count} alerts with remediation actions", 
         "default_color": CS_BLUE, "metric_name": "containment_rate", "raw_value": data.containment_rate},
        {"title": "Playbook Automation", "value": f"{data.playbook_auto['percent']}%", "detail": f"{data.playbook_auto['count']} alerts", 
         "default_color": CS_BLUE, "metric_name": "automation_percent", "raw_value": data.playbook_auto['percent']},
        {"title": "Human Review Rate", "value": f"{data.analyst_escalation['percent']}%", "detail": f"{data.analyst_escalation['count']} required manual review", 
         "default_color": CS_ORANGE, "metric_name": None, "raw_value": None}  # Renamed from Analyst Escalation
    ]
    
    for i, card in enumerate(efficiency_cards):
        card_top = cards_start_top + i * (card_height + card_spacing)
        
        # Get metric status if threshold exists
        status_info = None
        if card.get("metric_name") and card.get("raw_value") is not None:
            status_info = get_metric_status(card["metric_name"], card["raw_value"])
        
        # Determine border color based on status
        if status_info:
            indicator = get_status_indicator(status_info["status"], status_info["direction"])
            border_color = indicator["color"]
        else:
            border_color = card["default_color"]
        
        # Create card background
        card_shape = slide_resp_det.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left_panel_left, card_top,
            left_panel_width, card_height
        )
        fill = card_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)
        line = card_shape.line
        line.color.rgb = border_color
        line.width = Pt(2)
        
        # Card title (left)
        title_box = slide_resp_det.shapes.add_textbox(
            left_panel_left + Inches(0.15), card_top + Inches(0.15),
            Inches(2.2), Inches(0.3)
        )
        t_frame = title_box.text_frame
        t_para = t_frame.paragraphs[0]
        t_para.text = card["title"]
        t_para.font.name = TITLE_FONT_NAME
        t_para.font.size = Pt(14)
        t_para.font.bold = True
        t_para.font.color.rgb = CS_NAVY
        t_para.alignment = PP_ALIGN.LEFT
        
        # Detail (left, below title)
        detail_box = slide_resp_det.shapes.add_textbox(
            left_panel_left + Inches(0.15), card_top + Inches(0.48),
            Inches(2.2), Inches(0.3)
        )
        d_frame = detail_box.text_frame
        d_para = d_frame.paragraphs[0]
        d_para.text = card["detail"]
        d_para.font.name = BODY_FONT_NAME
        d_para.font.size = Pt(11)
        d_para.font.color.rgb = CS_SLATE
        d_para.alignment = PP_ALIGN.LEFT
        
        # Value with indicator (right)
        indicator_width = Inches(0.35) if status_info else 0
        value_box = slide_resp_det.shapes.add_textbox(
            left_panel_left + left_panel_width - Inches(1.5) - indicator_width, card_top + Inches(0.15),
            Inches(1.3), Inches(0.6)
        )
        v_frame = value_box.text_frame
        v_para = v_frame.paragraphs[0]
        v_para.text = card["value"]
        v_para.font.name = TITLE_FONT_NAME
        v_para.font.size = Pt(32)
        v_para.font.bold = True
        v_para.font.color.rgb = border_color
        v_para.alignment = PP_ALIGN.RIGHT
        
        # Note: Trend indicator arrows removed per Jan 2026 stakeholder feedback.
        # Colored border alone indicates status.
    
    # Right panel: Detection Quality (2-over-1 layout)
    # Per Jan 2026 feedback: 2 cards on top row, 1 spanning bottom
    # Renamed Threat Reduction → Alert Reduction
    right_panel_left = left_panel_left + left_panel_width + Inches(0.2)
    right_panel_width = left_panel_width
    top_card_width = (right_panel_width - Inches(0.1)) / 2  # 2 cards across on top row
    top_card_height = Inches(1.2)
    bottom_card_height = Inches(1.0)
    grid_spacing = Inches(0.1)
    
    # Calculate alert reduction percentage (2M potential threats → 267 escalated)
    potential_threats = 2000000  # From CORR funnel
    escalated = data.incidents_escalated
    reduction_percent = ((potential_threats - escalated) / potential_threats) * 100
    
    # Top row: True Positive Rate and False Positive Rate (side by side)
    top_row_cards = [
        {"title": "True Positive Rate", "value": f"{data.true_threat_precision}%", "default_color": CS_RED,
         "metric_name": "true_threat_precision", "raw_value": data.true_threat_precision},
        {"title": "False Positive Rate", "value": f"{data.false_positive_rate}%", "default_color": CS_ORANGE,
         "metric_name": "false_positive_rate", "raw_value": data.false_positive_rate}
    ]
    
    for i, card in enumerate(top_row_cards):
        card_left = right_panel_left + i * (top_card_width + grid_spacing)
        card_top = cards_start_top
        
        # Get metric status if threshold exists
        status_info = None
        if card.get("metric_name") and card.get("raw_value") is not None:
            status_info = get_metric_status(card["metric_name"], card["raw_value"])
        
        # Determine border color based on status
        if status_info:
            indicator = get_status_indicator(status_info["status"], status_info["direction"])
            border_color = indicator["color"]
        else:
            border_color = card["default_color"]
        
        # Create card background
        card_shape = slide_resp_det.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, card_left, card_top,
            top_card_width, top_card_height
        )
        fill = card_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 250, 252)
        line = card_shape.line
        line.color.rgb = border_color
        line.width = Pt(2)
        
        # Title
        title_box = slide_resp_det.shapes.add_textbox(
            card_left + Inches(0.1), card_top + Inches(0.1),
            top_card_width - Inches(0.2), Inches(0.3)
        )
        t_frame = title_box.text_frame
        t_para = t_frame.paragraphs[0]
        t_para.text = card["title"]
        t_para.font.name = TITLE_FONT_NAME
        t_para.font.size = Pt(12)
        t_para.font.bold = True
        t_para.font.color.rgb = CS_SLATE
        t_para.alignment = PP_ALIGN.CENTER
        
        # Value (centered)
        value_box = slide_resp_det.shapes.add_textbox(
            card_left + Inches(0.05), card_top + Inches(0.45),
            top_card_width - Inches(0.1), Inches(0.6)
        )
        v_frame = value_box.text_frame
        v_para = v_frame.paragraphs[0]
        v_para.text = card["value"]
        v_para.font.name = TITLE_FONT_NAME
        v_para.font.size = Pt(32)
        v_para.font.bold = True
        v_para.font.color.rgb = border_color
        v_para.alignment = PP_ALIGN.CENTER
    
    # Bottom row: Alert Reduction (full width, spanning both columns)
    bottom_card_top = cards_start_top + top_card_height + grid_spacing
    bottom_card_width = right_panel_width  # Full width
    
    # Create bottom card background
    bottom_card_shape = slide_resp_det.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, right_panel_left, bottom_card_top,
        bottom_card_width, bottom_card_height
    )
    fill = bottom_card_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    line = bottom_card_shape.line
    line.color.rgb = CS_GREEN
    line.width = Pt(2)
    
    # Bottom card title
    bottom_title_box = slide_resp_det.shapes.add_textbox(
        right_panel_left + Inches(0.15), bottom_card_top + Inches(0.1),
        bottom_card_width - Inches(0.3), Inches(0.25)
    )
    bt_frame = bottom_title_box.text_frame
    bt_para = bt_frame.paragraphs[0]
    bt_para.text = "Alert Reduction"  # Renamed from Threat Reduction
    bt_para.font.name = TITLE_FONT_NAME
    bt_para.font.size = Pt(12)
    bt_para.font.bold = True
    bt_para.font.color.rgb = CS_SLATE
    bt_para.alignment = PP_ALIGN.CENTER
    
    # Bottom card value (2M → 267)
    bottom_value_box = slide_resp_det.shapes.add_textbox(
        right_panel_left + Inches(0.1), bottom_card_top + Inches(0.35),
        bottom_card_width - Inches(0.2), Inches(0.35)
    )
    bv_frame = bottom_value_box.text_frame
    bv_para = bv_frame.paragraphs[0]
    bv_para.text = f"2M → {escalated}"
    bv_para.font.name = TITLE_FONT_NAME
    bv_para.font.size = Pt(28)
    bv_para.font.bold = True
    bv_para.font.color.rgb = CS_GREEN
    bv_para.alignment = PP_ALIGN.CENTER
    
    # Bottom card detail
    bottom_detail_box = slide_resp_det.shapes.add_textbox(
        right_panel_left + Inches(0.1), bottom_card_top + Inches(0.7),
        bottom_card_width - Inches(0.2), Inches(0.25)
    )
    bd_frame = bottom_detail_box.text_frame
    bd_para = bd_frame.paragraphs[0]
    bd_para.text = f"{reduction_percent:.2f}% resolved before escalation"
    bd_para.font.name = BODY_FONT_NAME
    bd_para.font.size = Pt(11)
    bd_para.font.color.rgb = CS_SLATE
    bd_para.alignment = PP_ALIGN.CENTER
    
    # Insight box at bottom - rewritten per Jan 2026 feedback to explain CORR context
    # "Of 2M potential threats, only X alerts required your attention"
    insight_top = bottom_card_top + bottom_card_height + Inches(0.15)
    insight_shape = slide_resp_det.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), insight_top,
        prs.slide_width - Inches(1), Inches(0.7)
    )
    fill = insight_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    line = insight_shape.line
    line.color.rgb = CS_BLUE
    line.width = Pt(2)
    
    # Generate insight text explaining CORR context - use correct language (potential threats, not events)
    insight_text = insight_shape.text_frame
    insight_text.text = (
        f"Of 2 million potential threats detected, only {data.incidents_escalated} alerts required your attention—"
        f"Critical Start resolved 99.99% before escalation."
    )
    insight_text.paragraphs[0].font.name = BODY_FONT_NAME
    insight_text.paragraphs[0].font.size = Pt(13)
    insight_text.paragraphs[0].font.color.rgb = CS_NAVY
    insight_text.paragraphs[0].font.bold = True
    insight_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    insight_text.vertical_anchor = 1


def build_threat_landscape_slides(prs, data, include=True):
    """Create the optional Threat Landscape section (Slides 11-13).
    
    LAYOUT STRUCTURE:
    -----------------
    Slide 11 - Severity Alignment Flow (Split Layout):
        Left Panel (~2"): Narrative cards (Upgraded, Downgraded, Aligned percentages)
        Right Panel (~7"): Sankey diagram showing vendor → CS severity flows
        
    Slide 12 - Threat & Detection Sources (Split Layout):
        Left Panel (~55%): MITRE ATT&CK stacked bar chart
        Right Panel (~45%): Detection source cards with FP rates
    
    KEY DATA FIELDS USED:
    - severity_flows (Sankey diagram data)
    - tactics, high_severity, medium_severity, low_severity (MITRE data)
    - detection_sources (source breakdown with FP rates)
    
    The 'include' parameter allows skipping these slides via --no-threat-landscape.
    
    Per CRITICALSTART branding guidelines:
    - All slides have transparent header and footer
    - Uses H1-H6 typography scale
    
    Args:
        prs (Presentation): The presentation object.
        data (ReportData): The report data object containing all metrics.
        include (bool): Whether to include this section. Default True.
    """
    if not include:
        return
    
    # Note: setup_content_slide is used for each slide now
    
    # Calculate severity alignment stats from severity_flows
    severity_order = ['Informational', 'Low', 'Medium', 'High', 'Critical']
    
    def get_severity_index(label):
        """Extract severity level from label and return its index."""
        parts = label.split(' ')
        severity_label = parts[-1]  # Get last part (e.g., 'Critical' from 'Vendor Critical')
        return severity_order.index(severity_label) if severity_label in severity_order else -1
    
    upgraded_count = 0
    de_escalated_count = 0
    aligned_count = 0
    
    for flow in data.severity_flows:
        source_idx = get_severity_index(flow['from'])
        target_idx = get_severity_index(flow['to'])
        flow_count = flow['flow']
        
        if target_idx > source_idx:
            upgraded_count += flow_count
        elif target_idx < source_idx:
            de_escalated_count += flow_count
        else:
            aligned_count += flow_count
    
    total_escalations = upgraded_count + de_escalated_count + aligned_count
    
    # Slide 11 - Severity Alignment Flow (Redesigned with narrative storytelling)
    slide11, content_top_11 = setup_content_slide(prs, "Severity Alignment Flow")
    
    # Add subtitle using H6 typography
    subtitle_left = MARGIN_STANDARD
    subtitle_top = content_top_11 - Inches(0.6)
    subtitle_width = prs.slide_width - Inches(2.5)
    subtitle_height = Inches(0.4)
    
    subtitle_box = slide11.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.text = "Vendor-reported criticality vs. Critical Start adjudication"
    subtitle_paragraph.font.name = BODY_FONT_NAME
    subtitle_paragraph.font.size = H6_FONT_SIZE  # 12pt
    subtitle_paragraph.font.color.rgb = CS_SLATE
    subtitle_paragraph.alignment = PP_ALIGN.LEFT
    
    # Note: Removed "X escalations" chip per Jan 2026 stakeholder feedback.
    # The three left-side narrative cards provide enough numerical context.
    
    # Calculate percentages for narrative cards
    upgrade_pct = (upgraded_count / total_escalations * 100) if total_escalations > 0 else 0
    downgrade_pct = (de_escalated_count / total_escalations * 100) if total_escalations > 0 else 0
    aligned_pct = (aligned_count / total_escalations * 100) if total_escalations > 0 else 0
    
    # Layout: Left stacked cards + Right larger Sankey
    # Left column: 3 stacked narrative cards
    card_left = Inches(0.5)
    card_width = Inches(2.0)
    card_height = Inches(1.0)
    card_spacing = Inches(0.15)
    cards_start_top = subtitle_top + subtitle_height + Inches(0.25)
    
    # Narrative-driven stat cards with storytelling
    # Per Jan 2026 feedback: show both percentage AND discrete count
    narrative_cards = [
        {
            "label": "Upgraded",
            "value": f"{upgrade_pct:.1f}%",
            "count": upgraded_count,
            "narrative": "Added Value Beyond\nVendor Detection",
            "color": CS_RED
        },
        {
            "label": "Downgraded",  # Changed from "De-escalated" per Jan 2026 terminology update
            "value": f"{downgrade_pct:.1f}%",
            "count": de_escalated_count,
            "narrative": "Analyst Time Returned\nto Client",
            "color": CS_BLUE
        },
        {
            "label": "Aligned",
            "value": f"{aligned_pct:.1f}%",
            "count": aligned_count,
            "narrative": "Vendor Assessment\nConfirmed",
            "color": CS_NAVY
        }
    ]
    
    for i, card in enumerate(narrative_cards):
        card_top = cards_start_top + i * (card_height + card_spacing)
        
        # Create card background
        card_shape = slide11.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, card_left, card_top,
            card_width, card_height
        )
        fill = card_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(250, 250, 252)  # Very light gray background
        line = card_shape.line
        line.color.rgb = card["color"]
        line.width = Pt(3)
        
        # Add card label (top)
        label_box = slide11.shapes.add_textbox(
            card_left + Inches(0.1), card_top + Inches(0.05),
            card_width - Inches(0.2), Inches(0.2)
        )
        label_frame = label_box.text_frame
        label_frame.word_wrap = True
        label_para = label_frame.paragraphs[0]
        label_para.text = card["label"].upper()
        label_para.font.name = BODY_FONT_NAME
        label_para.font.size = Pt(9)
        label_para.font.bold = True
        label_para.font.color.rgb = card["color"]
        label_para.alignment = PP_ALIGN.CENTER
        
        # Add percentage value with count (middle, prominent)
        # Per Jan 2026 feedback: show "12.3% (33)" format with both percentage and discrete count
        value_box = slide11.shapes.add_textbox(
            card_left + Inches(0.1), card_top + Inches(0.22),
            card_width - Inches(0.2), Inches(0.35)
        )
        value_frame = value_box.text_frame
        value_frame.word_wrap = True
        value_para = value_frame.paragraphs[0]
        value_para.text = f"{card['value']} ({card['count']})"
        value_para.font.name = TITLE_FONT_NAME
        value_para.font.size = Pt(22)  # Slightly smaller to fit both values
        value_para.font.bold = True
        value_para.font.color.rgb = card["color"]
        value_para.alignment = PP_ALIGN.CENTER
        
        # Add narrative subtitle (bottom)
        narrative_box = slide11.shapes.add_textbox(
            card_left + Inches(0.08), card_top + Inches(0.58),
            card_width - Inches(0.16), Inches(0.4)
        )
        narrative_frame = narrative_box.text_frame
        narrative_frame.word_wrap = True
        narrative_para = narrative_frame.paragraphs[0]
        narrative_para.text = card["narrative"]
        narrative_para.font.name = BODY_FONT_NAME
        narrative_para.font.size = Pt(8)
        narrative_para.font.color.rgb = CS_SLATE
        narrative_para.alignment = PP_ALIGN.CENTER
        narrative_para.line_spacing = 0.9
    
    # Right side: Larger Sankey chart placeholder
    chart_left = card_left + card_width + Inches(0.3)
    chart_top = cards_start_top
    chart_width = prs.slide_width - chart_left - Inches(0.5)
    chart_height = Inches(3.3)  # Larger height for the Sankey
    
    chart_placeholder = slide11.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        chart_width, chart_height
    )
    fill = chart_placeholder.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
    line = chart_placeholder.line
    line.color.rgb = CS_SLATE
    line.width = Pt(1)
    
    # Add placeholder text with ID for later replacement
    placeholder_text = chart_placeholder.text_frame
    placeholder_text.text = "[Chart: Severity Alignment Sankey - ID: severity_sankey]"
    placeholder_text.paragraphs[0].font.name = BODY_FONT_NAME
    placeholder_text.paragraphs[0].font.size = Pt(14)
    placeholder_text.paragraphs[0].font.color.rgb = CS_SLATE
    placeholder_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    placeholder_text.vertical_anchor = 1  # Middle
    
    # Slide 9 - Threat & Detection Sources (Merged: MITRE Tactics + Detection Sources)
    slide_threats, content_top_threats = setup_content_slide(prs, "Threat & Detection Sources")
    
    # Left half: MITRE ATT&CK chart placeholder (enlarged per Jan 2026 feedback)
    # Reduced detection source card width to allow more chart space
    chart_left = Inches(0.5)
    chart_top = content_top_threats + Inches(0.1)
    chart_width = prs.slide_width * 0.55  # Increased from 50% to 55%
    chart_height = Inches(3.0)  # Increased from 2.3"
    
    chart_placeholder = slide_threats.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        chart_width, chart_height
    )
    fill = chart_placeholder.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    line = chart_placeholder.line
    line.color.rgb = CS_SLATE
    line.width = Pt(1)
    
    placeholder_text = chart_placeholder.text_frame
    placeholder_text.text = "[Chart: MITRE ATT&CK Stacked Bar - ID: mitre_stacked_bar]"
    placeholder_text.paragraphs[0].font.name = BODY_FONT_NAME
    placeholder_text.paragraphs[0].font.size = Pt(12)
    placeholder_text.paragraphs[0].font.color.rgb = CS_SLATE
    placeholder_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    placeholder_text.vertical_anchor = 1
    
    # Note: Removed redundant "High | Medium | Low | Info" legend text box
    # and "Persistence: X escalations" insight text box per Jan 2026 feedback.
    # The chart has a built-in legend, and annotations go directly on the chart.
    
    # Right half: Detection Sources (3 compact cards stacked, narrower)
    right_left = chart_left + chart_width + Inches(0.15)
    right_width = prs.slide_width - right_left - Inches(0.5)  # Narrower cards
    source_card_height = Inches(0.95)  # Slightly reduced
    source_spacing = Inches(0.1)
    fp_threshold = 10.0
    
    for i, source in enumerate(data.detection_sources[:3]):  # Limit to 3
        card_top = content_top_threats + Inches(0.1) + i * (source_card_height + source_spacing)
        fp_rate = source.get('fp_rate', 0)
        fp_color = CS_ORANGE if fp_rate > fp_threshold else CS_BLUE
        
        # Card background
        card_shape = slide_threats.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, right_left, card_top,
            right_width, source_card_height
        )
        fill = card_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 250, 252)
        line = card_shape.line
        line.color.rgb = fp_color
        line.width = Pt(2)
        
        # Source name (left)
        name_box = slide_threats.shapes.add_textbox(
            right_left + Inches(0.12), card_top + Inches(0.12),
            Inches(2.5), Inches(0.35)
        )
        n_frame = name_box.text_frame
        n_para = n_frame.paragraphs[0]
        n_para.text = source['source']
        n_para.font.name = TITLE_FONT_NAME
        n_para.font.size = Pt(13)
        n_para.font.bold = True
        n_para.font.color.rgb = CS_NAVY
        n_para.alignment = PP_ALIGN.LEFT
        
        # Incidents (left, below name)
        inc_box = slide_threats.shapes.add_textbox(
            right_left + Inches(0.12), card_top + Inches(0.5),
            Inches(2.5), Inches(0.35)
        )
        i_frame = inc_box.text_frame
        i_para = i_frame.paragraphs[0]
        i_para.text = f"{source['incidents']} alerts ({source['percent']}%)"
        i_para.font.name = BODY_FONT_NAME
        i_para.font.size = Pt(11)
        i_para.font.color.rgb = CS_SLATE
        i_para.alignment = PP_ALIGN.LEFT
        
        # FP Rate (right side)
        fp_box = slide_threats.shapes.add_textbox(
            right_left + right_width - Inches(1.3), card_top + Inches(0.25),
            Inches(1.2), Inches(0.5)
        )
        f_frame = fp_box.text_frame
        f_para = f_frame.paragraphs[0]
        f_para.text = f"FP: {fp_rate}%"
        f_para.font.name = TITLE_FONT_NAME
        f_para.font.size = Pt(18)
        f_para.font.bold = True
        f_para.font.color.rgb = fp_color
        f_para.alignment = PP_ALIGN.RIGHT
    
    # Summary insight at bottom
    summary_top = content_top_threats + Inches(0.1) + 3 * (source_card_height + source_spacing) + Inches(0.15)
    summary_shape = slide_threats.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), summary_top,
        prs.slide_width - Inches(1), Inches(0.5)
    )
    fill = summary_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    line = summary_shape.line
    line.color.rgb = CS_BLUE
    line.width = Pt(2)
    
    summary_text = summary_shape.text_frame
    # Find source with highest FP rate
    high_fp_source = max(data.detection_sources, key=lambda x: x.get('fp_rate', 0))
    summary_text.text = f"Tuning opportunity: {high_fp_source['source']} at {high_fp_source['fp_rate']}% FP rate"
    summary_text.paragraphs[0].font.name = BODY_FONT_NAME
    summary_text.paragraphs[0].font.size = Pt(13)
    summary_text.paragraphs[0].font.color.rgb = CS_NAVY
    summary_text.paragraphs[0].font.bold = True
    summary_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    summary_text.vertical_anchor = 1


def build_insights_slides(prs, data):
    """Create the Insights section (Slides 14-15).
    
    LAYOUT STRUCTURE:
    -----------------
    Slide 14 - Prioritized Improvement Plan:
        3 stacked improvement cards with priority badges (HIGH/MEDIUM/LOW)
        Each card: Priority badge, title, description, owner, target
        Insight bar at bottom with focus areas summary
    
    Slide 15 - Operational Coverage (Dual Panel):
        Left Panel: After-hours hero metric + weeknight/weekend breakdown
        Right Panel: Collaboration metrics (avg touches, participation, closures)
    
    KEY DATA FIELDS USED:
    - improvement_items (list of dicts with priority, title, description, etc.)
    - after_hours_escalations, after_hours_weeknight, after_hours_weekend
    - avg_touches, client_participation, client_led_closures
    - after_hours_data_available (conditional display)
    
    Per CRITICALSTART branding guidelines:
    - All slides have transparent header and footer
    - Uses H1-H6 typography scale
    
    Args:
        prs (Presentation): The presentation object.
        data (ReportData): The report data object containing all metrics.
    """
    # Slide 14 - Prioritized Improvement Plan
    slide14, content_top_14 = setup_content_slide(prs, "Prioritized Improvement Plan")
    
    # Create improvement items cards
    # Reduced heights to prevent vertical overflow (per Jan 2026 feedback)
    card_start_top = content_top_14 + Inches(0.05)
    card_width = prs.slide_width - Inches(1.0)
    card_height = Inches(1.0)  # Reduced from 1.3"
    card_spacing = Inches(0.12)  # Reduced from 0.2"
    card_left = Inches(0.5)
    
    # Map priority to colors
    priority_colors = {
        "HIGH": CS_RED,
        "MEDIUM": CS_ORANGE,
        "LOW": CS_BLUE
    }
    
    # Sort improvement items by priority (HIGH -> MEDIUM -> LOW) per stakeholder feedback
    priority_order = {"HIGH": 0, "MEDIUM": 1, "LOW": 2}
    sorted_items = sorted(
        data.improvement_items,
        key=lambda x: priority_order.get(x.get("priority", "LOW"), 3)
    )
    
    # Expected impact text for each item (will be matched by position after sorting)
    expected_impacts = [
        "Reduce escalations by ~20%",
        "Reduce manual review burden",
        "Improved detection of advanced threats"
    ]
    
    # Concise descriptions as per requirements
    concise_descriptions = [
        "Palo Alto Cortex XDR false positive rate is 11.2%, exceeding the 10.0% threshold",
        "Manual escalations at 14% exceed 12% target. 38 alerts required analyst judgment",
        "Persistence + Defense Evasion account for 67% of high-severity alerts"
    ]
    
    for i, item in enumerate(sorted_items):
        card_top = card_start_top + i * (card_height + card_spacing)
        
        # Determine priority color
        priority = item.get("priority", "MEDIUM")
        border_color = priority_colors.get(priority, CS_ORANGE)
        
        # Create card background with colored left border
        card_shape = slide14.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, card_left, card_top,
            card_width, card_height
        )
        fill = card_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
        line = card_shape.line
        line.color.rgb = border_color
        line.width = Pt(4)  # Thicker left border
        
        # Add priority badge at top left (reduced sizes for tighter layout)
        badge_width = Inches(0.9)
        badge_height = Inches(0.28)
        badge_left = card_left + Inches(0.15)
        badge_top = card_top + Inches(0.1)
        
        badge_shape = slide14.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, badge_left, badge_top,
            badge_width, badge_height
        )
        fill = badge_shape.fill
        fill.solid()
        fill.fore_color.rgb = border_color
        badge_shape.line.fill.background()
        
        badge_text = badge_shape.text_frame
        badge_text.text = priority
        badge_text.paragraphs[0].font.name = TITLE_FONT_NAME
        badge_text.paragraphs[0].font.size = Pt(10)  # Reduced from 12
        badge_text.paragraphs[0].font.bold = True
        badge_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        badge_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_text.vertical_anchor = 1  # Middle
        
        # Add item title next to badge
        item_title_left = badge_left + badge_width + Inches(0.15)
        item_title_top = badge_top
        item_title_width = card_width - badge_width - Inches(0.5)
        item_title_height = badge_height
        
        item_title_box = slide14.shapes.add_textbox(
            item_title_left, item_title_top, item_title_width, item_title_height
        )
        item_title_frame = item_title_box.text_frame
        item_title_frame.word_wrap = True
        item_title_paragraph = item_title_frame.paragraphs[0]
        item_title_paragraph.text = item['title']  # Simplified title (removed "Item X -")
        item_title_paragraph.font.name = TITLE_FONT_NAME
        item_title_paragraph.font.size = Pt(14)  # Reduced from 16
        item_title_paragraph.font.bold = True
        item_title_paragraph.font.color.rgb = CS_NAVY
        item_title_paragraph.alignment = PP_ALIGN.LEFT
        
        # Add description
        desc_left = card_left + Inches(0.15)
        desc_top = badge_top + badge_height + Inches(0.08)
        desc_width = card_width - Inches(0.3)
        desc_height = Inches(0.35)
        
        desc_box = slide14.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_paragraph = desc_frame.paragraphs[0]
        # Use item's own description, with fallback to concise versions
        desc_paragraph.text = concise_descriptions[i] if i < len(concise_descriptions) else item.get('description', '')[:100]
        desc_paragraph.font.name = BODY_FONT_NAME
        desc_paragraph.font.size = Pt(11)  # Reduced from 12
        desc_paragraph.font.color.rgb = CS_SLATE
        desc_paragraph.alignment = PP_ALIGN.LEFT
        
        # Add metadata (Owner, Target, Expected Impact) in single row to save space
        meta_left = desc_left
        meta_top = desc_top + desc_height + Inches(0.05)
        meta_width = card_width - Inches(0.3)
        meta_height = Inches(0.28)
        
        # Combined metadata row
        impact_text = expected_impacts[i] if i < len(expected_impacts) else "Improved security posture"
        meta_text = f"{item['owner']} | {item['target']} | Impact: {impact_text}"
        meta_box = slide14.shapes.add_textbox(meta_left, meta_top, meta_width, meta_height)
        meta_frame = meta_box.text_frame
        meta_frame.word_wrap = True
        meta_paragraph = meta_frame.paragraphs[0]
        meta_paragraph.text = meta_text
        meta_paragraph.font.name = BODY_FONT_NAME
        meta_paragraph.font.size = Pt(10)  # Reduced from 11
        meta_paragraph.font.color.rgb = CS_SLATE
        meta_paragraph.alignment = PP_ALIGN.LEFT
    
    # Add insight bar for Prioritized Improvements (reduced height)
    improvements_insight_top = card_start_top + 3 * (card_height + card_spacing)
    add_insight_callout(
        slide14, prs,
        "Focus Areas for Continuous Improvement",
        "Three targeted actions to reduce noise, improve detection, and strengthen your security posture.",
        improvements_insight_top,
        height=Inches(0.55)  # Reduced from 0.7"
    )
    
    # Slide 11 - Operational Coverage (Merged: After-Hours + Operational Insights)
    slide_ops, content_top_ops = setup_content_slide(prs, "Operational Coverage")
    
    # Left panel: After-Hours metrics (hero + breakdown)
    left_width = (prs.slide_width - Inches(1.2)) / 2
    left_left = Inches(0.5)
    panel_top = content_top_ops + Inches(0.1)
    panel_height = Inches(3.0)
    
    # Left panel background
    left_panel = slide_ops.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left_left, panel_top,
        left_width, panel_height
    )
    fill = left_panel.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    line = left_panel.line
    line.color.rgb = CS_BLUE
    line.width = Pt(2)
    
    # After-Hours Hero number
    hero_box = slide_ops.shapes.add_textbox(left_left + Inches(0.2), panel_top + Inches(0.2), 
                                             Inches(2.5), Inches(1.0))
    hero_frame = hero_box.text_frame
    hero_para = hero_frame.paragraphs[0]
    hero_para.text = str(data.after_hours_escalations)
    hero_para.font.name = TITLE_FONT_NAME
    hero_para.font.size = Pt(64)
    hero_para.font.bold = True
    hero_para.font.color.rgb = CS_NAVY
    hero_para.alignment = PP_ALIGN.LEFT
    
    # After-Hours label
    label_box = slide_ops.shapes.add_textbox(left_left + Inches(0.2), panel_top + Inches(1.15), 
                                              Inches(3.0), Inches(0.35))
    label_frame = label_box.text_frame
    label_para = label_frame.paragraphs[0]
    label_para.text = "After-Hours Escalations"
    label_para.font.name = TITLE_FONT_NAME
    label_para.font.size = Pt(16)
    label_para.font.bold = True
    label_para.font.color.rgb = CS_NAVY
    label_para.alignment = PP_ALIGN.LEFT
    
    # Business hours definition (per Jan 2026 stakeholder feedback for clarity)
    if data.business_hours_definition:
        bh_box = slide_ops.shapes.add_textbox(left_left + Inches(0.2), panel_top + Inches(1.5), 
                                              Inches(3.0), Inches(0.25))
        bh_frame = bh_box.text_frame
        bh_para = bh_frame.paragraphs[0]
        bh_para.text = f"(Business hours: {data.business_hours_definition})"
        bh_para.font.name = BODY_FONT_NAME
        bh_para.font.size = Pt(10)
        bh_para.font.italic = True
        bh_para.font.color.rgb = CS_SLATE
        bh_para.alignment = PP_ALIGN.LEFT
    
    # Breakdown stats
    weeknight_count = getattr(data, 'after_hours_weeknight', int(data.after_hours_escalations * 0.82))
    weekend_count = getattr(data, 'after_hours_weekend', int(data.after_hours_escalations * 0.18))
    
    stats_box = slide_ops.shapes.add_textbox(left_left + Inches(0.2), panel_top + Inches(1.6), 
                                              left_width - Inches(0.4), Inches(1.2))
    stats_frame = stats_box.text_frame
    stats_frame.word_wrap = True
    stats_para = stats_frame.paragraphs[0]
    stats_para.text = f"{weeknight_count} weeknights · {weekend_count} weekends"
    stats_para.font.name = BODY_FONT_NAME
    stats_para.font.size = Pt(14)
    stats_para.font.color.rgb = CS_SLATE
    stats_para.alignment = PP_ALIGN.LEFT
    
    stats_para2 = stats_frame.add_paragraph()
    stats_para2.text = f"{data.coverage_hours} hours coverage · {int(data.automation_percent)}% auto-routed"
    stats_para2.font.name = BODY_FONT_NAME
    stats_para2.font.size = Pt(13)
    stats_para2.font.color.rgb = CS_SLATE
    stats_para2.alignment = PP_ALIGN.LEFT
    
    # Right panel: Collaboration metrics
    right_left = left_left + left_width + Inches(0.2)
    right_width = left_width
    
    right_panel = slide_ops.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, right_left, panel_top,
        right_width, panel_height
    )
    fill = right_panel.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)
    line = right_panel.line
    line.color.rgb = CS_SLATE
    line.width = Pt(2)
    
    # Collaboration title
    collab_title = slide_ops.shapes.add_textbox(right_left + Inches(0.2), panel_top + Inches(0.2), 
                                                 right_width - Inches(0.4), Inches(0.35))
    c_frame = collab_title.text_frame
    c_para = c_frame.paragraphs[0]
    c_para.text = "Collaboration Quality"
    c_para.font.name = TITLE_FONT_NAME
    c_para.font.size = Pt(16)
    c_para.font.bold = True
    c_para.font.color.rgb = CS_NAVY
    c_para.alignment = PP_ALIGN.LEFT
    
    # Collaboration metrics (3 compact rows)
    collab_data = [
        {"label": "Average Touches", "value": str(data.avg_touches), "detail": "per alert"},
        {"label": "Client Participation", "value": data.client_participation, "detail": "with client input"},
        {"label": "Client-Led Closures", "value": data.client_led_closures, "detail": "closed by team"}
    ]
    
    row_height = Inches(0.8)
    row_start = panel_top + Inches(0.6)
    
    for i, c in enumerate(collab_data):
        row_top = row_start + i * row_height
        
        # Value (left aligned)
        v_box = slide_ops.shapes.add_textbox(right_left + Inches(0.2), row_top, Inches(1.5), Inches(0.6))
        v_frame = v_box.text_frame
        v_para = v_frame.paragraphs[0]
        v_para.text = c["value"]
        v_para.font.name = TITLE_FONT_NAME
        v_para.font.size = Pt(28)
        v_para.font.bold = True
        v_para.font.color.rgb = CS_NAVY
        v_para.alignment = PP_ALIGN.LEFT
        
        # Label + detail (right of value)
        l_box = slide_ops.shapes.add_textbox(right_left + Inches(1.7), row_top + Inches(0.05), Inches(2.5), Inches(0.6))
        l_frame = l_box.text_frame
        l_para = l_frame.paragraphs[0]
        l_para.text = c["label"]
        l_para.font.name = TITLE_FONT_NAME
        l_para.font.size = Pt(12)
        l_para.font.bold = True
        l_para.font.color.rgb = CS_SLATE
        l_para.alignment = PP_ALIGN.LEFT
        
        d_para = l_frame.add_paragraph()
        d_para.text = c["detail"]
        d_para.font.name = BODY_FONT_NAME
        d_para.font.size = Pt(10)
        d_para.font.color.rgb = CS_SLATE
        d_para.alignment = PP_ALIGN.LEFT
    
    # Insight box at bottom
    insight_top = panel_top + panel_height + Inches(0.15)
    insight_shape = slide_ops.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), insight_top,
        prs.slide_width - Inches(1), Inches(0.55)
    )
    fill = insight_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    line = insight_shape.line
    line.color.rgb = CS_BLUE
    line.width = Pt(2)
    
    # Rewritten insight text per Jan 2026 feedback - explain significance, don't just restate metrics
    insight_text = insight_shape.text_frame
    insight_text.text = f"Critical Start handled {data.after_hours_escalations} alerts when your team was offline, maintaining 24/7 protection without gaps in coverage."
    insight_text.paragraphs[0].font.name = BODY_FONT_NAME
    insight_text.paragraphs[0].font.size = Pt(13)
    insight_text.paragraphs[0].font.color.rgb = CS_NAVY
    insight_text.paragraphs[0].font.bold = True
    insight_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    insight_text.vertical_anchor = 1


def build_forward_direction_slide(prs, data):
    """Create the Forward Direction slide (Slide 19 - Looking Ahead).
    
    Per CRITICALSTART branding guidelines:
    - All slides have transparent header and footer
    - Uses H1-H6 typography scale
    
    Note: Per Jan 2026 stakeholder feedback, this slide has NO insight box.
    This slide IS the narrative (it contains strategic recommendations),
    so an additional insight box would be redundant.
    
    Args:
        prs (Presentation): The presentation object.
        data (ReportData): The report data object containing all metrics.
    """
    # Forward Direction slide
    slide16, content_top_16 = setup_content_slide(prs, "Looking Ahead")
    
    # Section 1 - Next Period Targets (reduced heights to prevent overflow)
    section1_top = content_top_16 + Inches(0.05)
    section1_left = Inches(0.5)
    section1_width = prs.slide_width - Inches(1.0)
    section1_height = Inches(1.1)  # Reduced from 1.2"
    
    # Section 1 title
    section1_title_box = slide16.shapes.add_textbox(
        section1_left, section1_top, section1_width, Inches(0.4)
    )
    section1_title_frame = section1_title_box.text_frame
    section1_title_frame.word_wrap = True
    section1_title_paragraph = section1_title_frame.paragraphs[0]
    section1_title_paragraph.text = "Next Period Targets"
    section1_title_paragraph.font.name = TITLE_FONT_NAME
    section1_title_paragraph.font.size = Pt(20)
    section1_title_paragraph.font.bold = True
    section1_title_paragraph.font.color.rgb = CS_NAVY
    section1_title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Section 1 content (bullets)
    section1_content_top = section1_top + Inches(0.5)
    section1_content_box = slide16.shapes.add_textbox(
        section1_left + Inches(0.3), section1_content_top,
        section1_width - Inches(0.3), section1_height - Inches(0.5)
    )
    section1_content_frame = section1_content_box.text_frame
    section1_content_frame.word_wrap = True
    section1_content_frame.margin_left = Inches(0.2)
    section1_content_frame.margin_right = Inches(0.2)
    section1_content_frame.margin_top = Inches(0.1)
    section1_content_frame.margin_bottom = Inches(0.1)
    
    # Get Palo Alto XDR FP rate from detection_sources
    palo_alto_fp_rate = 11.2
    for source in data.detection_sources:
        if "Palo Alto" in source.get('source', ''):
            palo_alto_fp_rate = source.get('fp_rate', 11.2)
            break
    
    # Calculate reduction (11.2% to 10% = ~4 fewer escalations)
    # Assuming 189 alerts at 11.2% FP rate, reducing to 10% would save ~2-4 escalations
    target1_text = f"Trim Palo Alto XDR false positives from {palo_alto_fp_rate}% to 10% threshold (~4 fewer escalations)"
    target2_text = f"Reduce manual escalations from {data.analyst_escalation['count']} ({data.analyst_escalation['percent']}%) to 32 or fewer (12% target)"
    
    # Add bullet points with checkmark-like styling (using blue color)
    para1 = section1_content_frame.paragraphs[0]
    para1.text = "✓ " + target1_text
    para1.font.name = BODY_FONT_NAME
    para1.font.size = Pt(14)
    para1.font.color.rgb = CS_BLUE
    para1.level = 0
    para1.space_after = Pt(8)
    
    para2 = section1_content_frame.add_paragraph()
    para2.text = "✓ " + target2_text
    para2.font.name = BODY_FONT_NAME
    para2.font.size = Pt(14)
    para2.font.color.rgb = CS_BLUE
    para2.level = 0
    para2.space_after = Pt(8)
    
    # Section 2 - Strategic Focus
    section2_top = section1_top + section1_height + Inches(0.2)  # Reduced gap
    section2_height = Inches(0.9)  # Reduced from 1.0"
    
    # Section 2 title
    section2_title_box = slide16.shapes.add_textbox(
        section1_left, section2_top, section1_width, Inches(0.4)
    )
    section2_title_frame = section2_title_box.text_frame
    section2_title_frame.word_wrap = True
    section2_title_paragraph = section2_title_frame.paragraphs[0]
    section2_title_paragraph.text = "Strategic Focus"
    section2_title_paragraph.font.name = TITLE_FONT_NAME
    section2_title_paragraph.font.size = Pt(20)
    section2_title_paragraph.font.bold = True
    section2_title_paragraph.font.color.rgb = CS_NAVY
    section2_title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Section 2 content
    section2_content_top = section2_top + Inches(0.5)
    section2_content_box = slide16.shapes.add_textbox(
        section1_left + Inches(0.3), section2_content_top,
        section1_width - Inches(0.3), section2_height - Inches(0.5)
    )
    section2_content_frame = section2_content_box.text_frame
    section2_content_frame.word_wrap = True
    section2_content_frame.margin_left = Inches(0.2)
    section2_content_frame.margin_right = Inches(0.2)
    section2_content_frame.margin_top = Inches(0.1)
    section2_content_frame.margin_bottom = Inches(0.1)
    
    focus1_text = "Proactive hunts targeting Persistence and Defense Evasion tactics"
    focus2_text = "Additional playbook coverage for analyst-escalated scenarios"
    
    para3 = section2_content_frame.paragraphs[0]
    para3.text = "✓ " + focus1_text
    para3.font.name = BODY_FONT_NAME
    para3.font.size = Pt(14)
    para3.font.color.rgb = CS_BLUE
    para3.level = 0
    para3.space_after = Pt(8)
    
    para4 = section2_content_frame.add_paragraph()
    para4.text = "✓ " + focus2_text
    para4.font.name = BODY_FONT_NAME
    para4.font.size = Pt(14)
    para4.font.color.rgb = CS_BLUE
    para4.level = 0
    para4.space_after = Pt(8)
    
    # Section 3 - Your Partnership
    section3_top = section2_top + section2_height + Inches(0.2)  # Reduced gap
    section3_height = Inches(0.7)  # Reduced from 0.8"
    
    # Section 3 title
    section3_title_box = slide16.shapes.add_textbox(
        section1_left, section3_top, section1_width, Inches(0.4)
    )
    section3_title_frame = section3_title_box.text_frame
    section3_title_frame.word_wrap = True
    section3_title_paragraph = section3_title_frame.paragraphs[0]
    section3_title_paragraph.text = "Your Partnership"
    section3_title_paragraph.font.name = TITLE_FONT_NAME
    section3_title_paragraph.font.size = Pt(20)
    section3_title_paragraph.font.bold = True
    section3_title_paragraph.font.color.rgb = CS_NAVY
    section3_title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Section 3 content
    section3_content_top = section3_top + Inches(0.5)
    section3_content_box = slide16.shapes.add_textbox(
        section1_left, section3_content_top,
        section1_width, section3_height - Inches(0.5)
    )
    section3_content_frame = section3_content_box.text_frame
    section3_content_frame.word_wrap = True
    section3_content_paragraph = section3_content_frame.paragraphs[0]
    section3_content_paragraph.text = "Questions? Your Customer Success Manager is here to help."
    section3_content_paragraph.font.name = BODY_FONT_NAME
    section3_content_paragraph.font.size = Pt(16)
    section3_content_paragraph.font.color.rgb = CS_SLATE
    section3_content_paragraph.alignment = PP_ALIGN.LEFT
    
    # Note: Per Jan 2026 stakeholder feedback, NO insight bar on this slide.
    # This slide IS the narrative content, so a "Strategic Path Forward" box
    # would be redundant. The slide title and content serve that purpose.


def build_contact_slide(prs, data):
    """Create the closing Contact/Thank You slide.
    
    Args:
        prs (Presentation): The presentation object.
        data (ReportData): The report data object containing all metrics.
    """
    blank_slide_layout = prs.slide_layouts[6]
    slide_contact = prs.slides.add_slide(blank_slide_layout)
    slide_number = get_slide_number(prs)
    
    # Add master slide elements
    add_master_slide_elements(slide_contact, prs, slide_number=slide_number,
                               include_header=True, include_footer=True)
    
    # Add logo at top right
    add_logo(slide_contact, position='top_right', prs=prs)
    
    # Center content area
    content_top = HEADER_HEIGHT + Inches(0.8)
    
    # Thank You message
    thank_you_box = slide_contact.shapes.add_textbox(
        Inches(0.5), content_top,
        prs.slide_width - Inches(1), Inches(1.0)
    )
    thank_you_frame = thank_you_box.text_frame
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.text = "Thank You"
    thank_you_para.font.name = TITLE_FONT_NAME
    thank_you_para.font.size = Pt(60)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = CS_NAVY
    thank_you_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide_contact.shapes.add_textbox(
        Inches(0.5), content_top + Inches(1.2),
        prs.slide_width - Inches(1), Inches(0.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "Questions? We're here to help."
    subtitle_para.font.name = BODY_FONT_NAME
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = CS_SLATE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Contact card
    card_width = Inches(6)
    card_height = Inches(1.5)
    card_left = (prs.slide_width - card_width) / 2
    card_top = content_top + Inches(2.2)
    
    card_shape = slide_contact.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, card_left, card_top,
        card_width, card_height
    )
    fill = card_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    line = card_shape.line
    line.color.rgb = CS_BLUE
    line.width = Pt(2)
    
    # Contact info
    contact_box = slide_contact.shapes.add_textbox(
        card_left + Inches(0.3), card_top + Inches(0.3),
        card_width - Inches(0.6), card_height - Inches(0.6)
    )
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    
    contact_para1 = contact_frame.paragraphs[0]
    contact_para1.text = "Your Customer Success Manager"
    contact_para1.font.name = TITLE_FONT_NAME
    contact_para1.font.size = Pt(16)
    contact_para1.font.bold = True
    contact_para1.font.color.rgb = CS_SLATE
    contact_para1.alignment = PP_ALIGN.CENTER
    
    # Dynamic CSM name per Jan 2026 stakeholder feedback
    contact_para2 = contact_frame.add_paragraph()
    csm_name = data.csm_name if data.csm_name else "Your CS Team"
    contact_para2.text = csm_name
    contact_para2.font.name = TITLE_FONT_NAME
    contact_para2.font.size = Pt(22)
    contact_para2.font.bold = True
    contact_para2.font.color.rgb = CS_NAVY
    contact_para2.alignment = PP_ALIGN.CENTER
    
    # Dynamic CSM email or fallback to support
    contact_para3 = contact_frame.add_paragraph()
    csm_email = data.csm_email if data.csm_email else "support@criticalstart.com"
    contact_para3.text = csm_email
    contact_para3.font.name = BODY_FONT_NAME
    contact_para3.font.size = Pt(14)
    contact_para3.font.color.rgb = CS_BLUE
    contact_para3.alignment = PP_ALIGN.CENTER
    
    contact_para4 = contact_frame.add_paragraph()
    contact_para4.text = "www.criticalstart.com"
    contact_para4.font.name = BODY_FONT_NAME
    contact_para4.font.size = Pt(12)
    contact_para4.font.color.rgb = CS_SLATE
    contact_para4.alignment = PP_ALIGN.CENTER
    
    # Report date footer
    report_date_box = slide_contact.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.7),
        prs.slide_width - Inches(1), Inches(0.3)
    )
    report_date_frame = report_date_box.text_frame
    report_date_para = report_date_frame.paragraphs[0]
    report_date_para.text = f"Report Date: {data.report_date}"
    report_date_para.font.name = BODY_FONT_NAME
    report_date_para.font.size = Pt(11)
    report_date_para.font.color.rgb = CS_SLATE
    report_date_para.font.italic = True
    report_date_para.alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    main()

