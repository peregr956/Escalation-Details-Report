"""Reusable layout helpers for PowerPoint presentation generation.

This module contains helper functions for creating common slide elements
like metric cards, insight callouts, data tables, and slide layouts.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from datetime import datetime

from constants import (
    CS_BLUE, CS_NAVY, CS_SLATE, CS_VIOLET, CS_RED, CS_ORANGE, CS_GREEN,
    TITLE_FONT_NAME, BODY_FONT_NAME,
    H1_FONT_SIZE, H2_FONT_SIZE, H3_FONT_SIZE, H4_FONT_SIZE,
    H5_FONT_SIZE, H6_FONT_SIZE, PARAGRAPH_FONT_SIZE, FOOTER_FONT_SIZE,
    TITLE_FONT_SIZE, BODY_FONT_SIZE,
    MARGIN_STANDARD, MARGIN_CONTENT, CARD_SPACING, HEADER_HEIGHT, FOOTER_HEIGHT,
    PRESENTATION_TITLE, PRESENTATION_INTENT, COPYRIGHT_TEXT,
    METRIC_THRESHOLDS
)


# =============================================================================
# PRESENTATION INITIALIZATION
# =============================================================================

def create_presentation():
    """Initialize 16:9 widescreen presentation.
    
    Returns:
        Presentation: A new Presentation object with 16:9 aspect ratio.
    """
    prs = Presentation()
    # Set slide size to 16:9 widescreen (10 x 5.625 inches)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs


def apply_branding(prs):
    """Set up slide masters with brand colors.
    
    Args:
        prs (Presentation): The presentation to apply branding to.
    """
    # Get the slide master
    slide_master = prs.slide_master
    
    # Set background color to white (default)
    background = slide_master.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def get_slide_number(prs):
    """Return the current slide count for numbering.
    
    Args:
        prs (Presentation): The presentation object.
    
    Returns:
        int: Number of slides currently in the presentation.
    """
    return len(prs.slides)


# =============================================================================
# MASTER SLIDE ELEMENTS
# =============================================================================

def add_master_slide_header(slide, prs, slide_number=None, include_header=True, text_color=None):
    """Add master slide header elements per CRITICALSTART branding guidelines.
    
    Header contains (all caps, transparent background):
    - Short Presentation Title (e.g., "ESCALATION REPORT")
    - "CRITICAL START"
    - Slide number
    
    Args:
        slide: The slide object to add header to.
        prs (Presentation): The presentation object for dimensions.
        slide_number: The slide number to display (optional).
        include_header: Whether to include header elements (False for title slides).
        text_color: Optional RGBColor for header text (defaults to CS_SLATE).
    """
    if not include_header:
        return
    
    if text_color is None:
        text_color = CS_SLATE
    
    header_top = Inches(0.15)
    header_height = HEADER_HEIGHT
    
    # Calculate positions for 3-column header
    col_width = (prs.slide_width - MARGIN_STANDARD * 2) / 3
    left_col = MARGIN_STANDARD
    center_col = MARGIN_STANDARD + col_width
    right_col = MARGIN_STANDARD + col_width * 2
    
    # Left: Presentation Title
    title_box = slide.shapes.add_textbox(left_col, header_top, col_width, header_height)
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = PRESENTATION_TITLE
    title_para.font.name = TITLE_FONT_NAME
    title_para.font.size = FOOTER_FONT_SIZE
    title_para.font.bold = False
    title_para.font.color.rgb = text_color
    title_para.alignment = PP_ALIGN.LEFT
    
    # Center: CRITICAL START
    cs_box = slide.shapes.add_textbox(center_col, header_top, col_width, header_height)
    cs_frame = cs_box.text_frame
    cs_para = cs_frame.paragraphs[0]
    cs_para.text = "CRITICAL START"
    cs_para.font.name = TITLE_FONT_NAME
    cs_para.font.size = FOOTER_FONT_SIZE
    cs_para.font.bold = False
    cs_para.font.color.rgb = text_color
    cs_para.alignment = PP_ALIGN.CENTER
    
    # Right: Slide number
    if slide_number is not None:
        num_box = slide.shapes.add_textbox(right_col, header_top, col_width, header_height)
        num_frame = num_box.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = str(slide_number)
        num_para.font.name = TITLE_FONT_NAME
        num_para.font.size = FOOTER_FONT_SIZE
        num_para.font.bold = False
        num_para.font.color.rgb = text_color
        num_para.alignment = PP_ALIGN.RIGHT


def add_master_slide_footer(slide, prs, date_text=None, include_footer=True, text_color=None):
    """Add master slide footer elements per CRITICALSTART branding guidelines.
    
    Footer contains (all caps, transparent background):
    - Date (e.g., "DECEMBER 2025")
    - "©2025 CRITICAL START"
    - Intent summary (e.g., "EBR")
    
    Args:
        slide: The slide object to add footer to.
        prs (Presentation): The presentation object for dimensions.
        date_text: Custom date text (defaults to current month/year).
        include_footer: Whether to include footer elements.
        text_color: Optional RGBColor for footer text (defaults to CS_SLATE).
    """
    if not include_footer:
        return
    
    if text_color is None:
        text_color = CS_SLATE
    
    footer_top = prs.slide_height - FOOTER_HEIGHT - Inches(0.1)
    footer_height = FOOTER_HEIGHT
    
    # Calculate positions for 3-column footer
    col_width = (prs.slide_width - MARGIN_STANDARD * 2) / 3
    left_col = MARGIN_STANDARD
    center_col = MARGIN_STANDARD + col_width
    right_col = MARGIN_STANDARD + col_width * 2
    
    # Default date if not provided
    if date_text is None:
        date_text = datetime.now().strftime("%B %Y").upper()
    
    # Left: Date
    date_box = slide.shapes.add_textbox(left_col, footer_top, col_width, footer_height)
    date_frame = date_box.text_frame
    date_para = date_frame.paragraphs[0]
    date_para.text = date_text
    date_para.font.name = TITLE_FONT_NAME
    date_para.font.size = FOOTER_FONT_SIZE
    date_para.font.bold = False
    date_para.font.color.rgb = text_color
    date_para.alignment = PP_ALIGN.LEFT
    
    # Center: Copyright
    copyright_box = slide.shapes.add_textbox(center_col, footer_top, col_width, footer_height)
    copyright_frame = copyright_box.text_frame
    copyright_para = copyright_frame.paragraphs[0]
    copyright_para.text = COPYRIGHT_TEXT
    copyright_para.font.name = TITLE_FONT_NAME
    copyright_para.font.size = FOOTER_FONT_SIZE
    copyright_para.font.bold = False
    copyright_para.font.color.rgb = text_color
    copyright_para.alignment = PP_ALIGN.CENTER
    
    # Right: Intent
    intent_box = slide.shapes.add_textbox(right_col, footer_top, col_width, footer_height)
    intent_frame = intent_box.text_frame
    intent_para = intent_frame.paragraphs[0]
    intent_para.text = PRESENTATION_INTENT
    intent_para.font.name = TITLE_FONT_NAME
    intent_para.font.size = FOOTER_FONT_SIZE
    intent_para.font.bold = False
    intent_para.font.color.rgb = text_color
    intent_para.alignment = PP_ALIGN.RIGHT


def add_master_slide_elements(slide, prs, slide_number=None,
                               include_header=True, include_footer=True,
                               date_text=None, text_color=None):
    """Add all master slide elements (header and footer) per branding guidelines.
    
    Args:
        slide: The slide object to add elements to.
        prs (Presentation): The presentation object for dimensions.
        slide_number: The slide number to display in header (optional).
        include_header: Whether to include header elements (False for title slides).
        include_footer: Whether to include footer elements.
        date_text: Custom date text for footer (defaults to current month/year).
        text_color: Optional RGBColor for header/footer text (defaults to CS_SLATE).
    """
    add_master_slide_header(slide, prs, slide_number, include_header, text_color)
    add_master_slide_footer(slide, prs, date_text, include_footer, text_color)


# =============================================================================
# LOGO HANDLING
# =============================================================================

def add_logo(slide, position='top_right', prs=None):
    """Add Critical Start logo to slides.
    
    Uses the approved PNG logo file (assets/Critical-Start-Stacked-Logo_0-2.png).
    
    Args:
        slide: The slide object to add the logo to.
        position (str): Position for logo placement. Options:
            - 'top_left': Top left corner
            - 'top_right': Top right corner (default)
            - 'bottom_left': Bottom left corner
            - 'bottom_right': Bottom right corner
        prs (Presentation, optional): Presentation object for getting dimensions.
    
    Returns:
        Picture: The picture shape object, or None if logo file not found.
    """
    # Use approved PNG logo
    logo_path = Path("assets/Critical-Start-Stacked-Logo_0-2.png")
    
    if not logo_path.exists():
        print(f"Warning: Logo file not found.")
        print(f"  Expected: assets/Critical-Start-Stacked-Logo_0-2.png")
        return None
    
    # Get original image dimensions to preserve aspect ratio
    from PIL import Image
    with Image.open(logo_path) as img:
        orig_width, orig_height = img.size
        aspect_ratio = orig_width / orig_height
    
    # Define logo width - height is calculated to preserve aspect ratio
    logo_width = Inches(0.7)
    logo_height = logo_width / aspect_ratio
    
    # Get slide dimensions
    if prs:
        slide_width = prs.slide_width
        slide_height = prs.slide_height
    else:
        slide_width = Inches(10)
        slide_height = Inches(5.625)
    
    if position == 'top_left':
        left = Inches(0.5)
        top = Inches(0.3)
    elif position == 'top_right':
        left = slide_width - logo_width - Inches(0.5)
        top = Inches(0.3)
    elif position == 'bottom_left':
        left = Inches(0.5)
        top = slide_height - logo_height - Inches(0.3)
    elif position == 'bottom_right':
        left = slide_width - logo_width - Inches(0.5)
        top = slide_height - logo_height - Inches(0.3)
    else:
        left = slide_width - logo_width - Inches(0.5)
        top = Inches(0.3)
    
    try:
        # Only specify width to let python-pptx preserve the original aspect ratio
        picture = slide.shapes.add_picture(str(logo_path), left, top, width=logo_width)
        return picture
    except Exception as e:
        print(f"Error adding logo: {e}")
        return None


# =============================================================================
# CONTENT SLIDE HELPERS
# =============================================================================

def setup_content_slide(prs, title_text, include_title=True):
    """Create a new content slide with standard layout per branding guidelines.
    
    This helper function creates a slide with:
    - Transparent header with title, CRITICAL START, slide number
    - Transparent footer with date, copyright, intent
    - Logo at top right
    - Optional slide title using H3 typography
    
    Args:
        prs (Presentation): The presentation object.
        title_text (str): The slide title text.
        include_title: Whether to add a slide title (default True).
    
    Returns:
        tuple: (slide, content_top) where content_top is the Y position for content.
    """
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide_number = get_slide_number(prs)
    
    # Add master slide elements
    add_master_slide_elements(slide, prs, slide_number=slide_number,
                               include_header=True, include_footer=True)
    
    # Add logo
    add_logo(slide, position='top_right', prs=prs)
    
    # Content starts below header
    content_top = HEADER_HEIGHT + Inches(0.3)
    
    # Add slide title using H3 typography
    if include_title and title_text:
        title_box = slide.shapes.add_textbox(
            MARGIN_STANDARD, content_top,
            prs.slide_width - Inches(2), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title_text
        title_para.font.name = TITLE_FONT_NAME
        title_para.font.size = H3_FONT_SIZE
        title_para.font.bold = True
        title_para.font.color.rgb = CS_NAVY
        title_para.alignment = PP_ALIGN.LEFT
        
        content_top = content_top + Inches(0.8)
    
    return slide, content_top


def create_gradient_background(prs, slide, gradient_type='blue_sweep'):
    """Create a gradient background for a slide using PowerPoint's native gradient fill.
    
    Args:
        prs (Presentation): The presentation object for dimensions.
        slide: The slide object to add the gradient to.
        gradient_type (str): Type of gradient to create. Options:
            - 'blue_sweep': Blue to Navy gradient (#009CDE → #004C97)
            - 'navy_solid': Solid navy background
    
    Returns:
        The background shape object.
    """
    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, prs.slide_height
    )
    background_shape.line.fill.background()
    
    if gradient_type == 'navy_solid':
        fill = background_shape.fill
        fill.solid()
        fill.fore_color.rgb = CS_NAVY
        return background_shape
    
    # For blue_sweep gradient
    fill = background_shape.fill
    fill.gradient()
    fill.gradient_angle = 0  # Left to right
    fill.gradient_stops[0].color.rgb = CS_BLUE
    fill.gradient_stops[1].color.rgb = CS_NAVY
    
    return background_shape


def add_header_bar(slide, prs, title_text, use_white_text=True):
    """Add a standardized header bar with title to a slide.
    
    Args:
        slide: The slide object to add the header to.
        prs (Presentation): The presentation object for dimensions.
        title_text (str): The title text to display.
        use_white_text (bool): Whether to use white text (for dark backgrounds).
    
    Returns:
        Tuple of (header_shape, title_textbox)
    """
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, HEADER_HEIGHT
    )
    fill = header_shape.fill
    fill.solid()
    fill.fore_color.rgb = CS_NAVY
    header_shape.line.fill.background()
    
    title_left = MARGIN_STANDARD
    title_top = Inches(0.1)
    title_width = prs.slide_width - Inches(2.5)
    title_height = Inches(0.6)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title_text
    title_paragraph.font.name = TITLE_FONT_NAME
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255) if use_white_text else CS_NAVY
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    return header_shape, title_box


# =============================================================================
# COMPONENT HELPERS
# =============================================================================

def add_insight_callout(slide, prs, title_text, body_text, top, height=Inches(1.2)):
    """Add a standardized 'What This Means' insight callout box.
    
    Args:
        slide: The slide object to add the callout to.
        prs (Presentation): The presentation object for dimensions.
        title_text (str): The callout title.
        body_text (str): The callout body text.
        top: The top position of the callout.
        height: The height of the callout box.
    
    Returns:
        The callout shape object.
    """
    callout_left = MARGIN_STANDARD
    callout_width = prs.slide_width - MARGIN_STANDARD * 2
    
    callout_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, callout_left, top,
        callout_width, height
    )
    fill = callout_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    line = callout_shape.line
    line.color.rgb = CS_BLUE
    line.width = Pt(3)
    
    callout_text = callout_shape.text_frame
    callout_text.margin_left = Inches(0.2)
    callout_text.margin_right = Inches(0.2)
    callout_text.margin_top = Inches(0.15)
    callout_text.margin_bottom = Inches(0.15)
    
    title_para = callout_text.paragraphs[0]
    title_para.text = title_text
    title_para.font.name = TITLE_FONT_NAME
    title_para.font.size = Pt(16)
    title_para.font.bold = True
    title_para.font.color.rgb = CS_NAVY
    title_para.alignment = PP_ALIGN.LEFT
    
    body_para = callout_text.add_paragraph()
    body_para.text = body_text
    body_para.font.name = BODY_FONT_NAME
    body_para.font.size = Pt(14)
    body_para.font.color.rgb = CS_SLATE
    body_para.alignment = PP_ALIGN.LEFT
    
    return callout_shape


def create_metric_card(slide, left, top, width, height,
                       value, label, context=None,
                       border_color=None, value_size=Pt(42)):
    """Create a standardized metric display card.
    
    Args:
        slide: The slide object to add the card to.
        left: Left position of the card.
        top: Top position of the card.
        width: Width of the card.
        height: Height of the card.
        value: The metric value to display (string).
        label: The label for the metric.
        context: Optional context text below the label.
        border_color: Optional border color (defaults to CS_BLUE).
        value_size: Font size for the value (default Pt(42)).
    
    Returns:
        The card shape object.
    """
    if border_color is None:
        border_color = CS_BLUE
    
    # Card background
    card_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = card_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    line = card_shape.line
    line.color.rgb = border_color
    line.width = Pt(3)
    
    # Value
    value_height = height * 0.45
    value_box = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.15),
        width - Inches(0.2), value_height
    )
    value_frame = value_box.text_frame
    value_para = value_frame.paragraphs[0]
    value_para.text = str(value)
    value_para.font.name = TITLE_FONT_NAME
    value_para.font.size = value_size
    value_para.font.bold = True
    value_para.font.color.rgb = CS_NAVY
    value_para.alignment = PP_ALIGN.CENTER
    
    # Label
    label_top = top + value_height + Inches(0.1)
    label_box = slide.shapes.add_textbox(
        left + Inches(0.1), label_top,
        width - Inches(0.2), Inches(0.4)
    )
    label_frame = label_box.text_frame
    label_para = label_frame.paragraphs[0]
    label_para.text = label
    label_para.font.name = TITLE_FONT_NAME
    label_para.font.size = Pt(14)
    label_para.font.bold = True
    label_para.font.color.rgb = border_color
    label_para.alignment = PP_ALIGN.CENTER
    
    # Context (if provided)
    if context:
        context_top = label_top + Inches(0.35)
        context_box = slide.shapes.add_textbox(
            left + Inches(0.1), context_top,
            width - Inches(0.2), Inches(0.4)
        )
        context_frame = context_box.text_frame
        context_para = context_frame.paragraphs[0]
        context_para.text = context
        context_para.font.name = BODY_FONT_NAME
        context_para.font.size = Pt(11)
        context_para.font.color.rgb = CS_SLATE
        context_para.alignment = PP_ALIGN.CENTER
    
    return card_shape


def create_data_table(slide, left, top, width, height,
                      headers, rows, header_color=None):
    """Create a standardized data table.
    
    Args:
        slide: The slide object to add the table to.
        left: Left position of the table.
        top: Top position of the table.
        width: Width of the table.
        height: Height of the table.
        headers: List of column header strings.
        rows: List of row data (each row is a list of values).
        header_color: Optional header background color (defaults to CS_NAVY).
    
    Returns:
        The table shape object.
    """
    if header_color is None:
        header_color = CS_NAVY
    
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, height
    ).table
    
    # Set column widths (equal distribution)
    col_width = width / num_cols
    for col_idx in range(num_cols):
        table.columns[col_idx].width = col_width
    
    # Style header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.name = TITLE_FONT_NAME
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Style data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 248, 248)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = BODY_FONT_NAME
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = CS_SLATE
            paragraph.alignment = PP_ALIGN.CENTER
    
    return table


# =============================================================================
# METRIC STATUS & TREND INDICATORS
# =============================================================================

def get_metric_status(metric_name, value):
    """Evaluate a metric value against defined thresholds.
    
    Args:
        metric_name (str): The metric identifier (must match key in METRIC_THRESHOLDS).
        value (float|int): The current value of the metric.
    
    Returns:
        dict: Contains 'status' ("good", "warning", or "bad") and 
              'direction' ("lower_is_better" or "higher_is_better").
              Returns None if metric_name not found in thresholds.
    
    Example:
        >>> get_metric_status("mttr", 126)
        {'status': 'good', 'direction': 'lower_is_better'}
        >>> get_metric_status("containment_rate", 85)
        {'status': 'bad', 'direction': 'higher_is_better'}
    """
    if metric_name not in METRIC_THRESHOLDS:
        return None
    
    config = METRIC_THRESHOLDS[metric_name]
    direction = config["direction"]
    good_threshold = config["good_threshold"]
    warning_threshold = config["warning_threshold"]
    
    if direction == "lower_is_better":
        # Lower values are better: good <= good_threshold, warning between, bad > warning_threshold
        if value <= good_threshold:
            status = "good"
        elif value <= warning_threshold:
            status = "warning"
        else:
            status = "bad"
    else:  # higher_is_better
        # Higher values are better: good >= good_threshold, warning between, bad < warning_threshold
        if value >= good_threshold:
            status = "good"
        elif value >= warning_threshold:
            status = "warning"
        else:
            status = "bad"
    
    return {
        "status": status,
        "direction": direction
    }


def get_status_indicator(status, direction):
    """Get the appropriate arrow character and color based on status and direction.
    
    Uses outlined-style Unicode arrows per FA6 brand guidelines:
    - ↑ (U+2191) for upward/elevated
    - ↓ (U+2193) for downward/depreciated
    - → (U+2192) for neutral/warning
    
    The arrow direction reflects whether the value is elevated (↑) or depreciated (↓),
    while the color reflects whether that's good or bad.
    
    Args:
        status (str): One of "good", "warning", or "bad".
        direction (str): One of "lower_is_better" or "higher_is_better".
    
    Returns:
        dict: Contains 'arrow' (Unicode character) and 'color' (RGBColor).
    
    Logic:
        - higher_is_better + good = ↑ green (high is good)
        - higher_is_better + bad = ↓ red (low is bad)
        - lower_is_better + good = ↓ green (low is good)
        - lower_is_better + bad = ↑ red (high is bad)
        - warning = → orange (neutral)
    """
    if status == "warning":
        return {"arrow": "→", "color": CS_ORANGE}
    
    if direction == "higher_is_better":
        if status == "good":
            return {"arrow": "↑", "color": CS_GREEN}
        else:  # bad
            return {"arrow": "↓", "color": CS_RED}
    else:  # lower_is_better
        if status == "good":
            return {"arrow": "↓", "color": CS_GREEN}
        else:  # bad
            return {"arrow": "↑", "color": CS_RED}


def add_trend_indicator(slide, left, top, status, direction, size=Pt(18)):
    """Add a trend indicator arrow to a slide.
    
    Args:
        slide: The slide object to add the indicator to.
        left: Left position of the indicator.
        top: Top position of the indicator.
        status (str): One of "good", "warning", or "bad".
        direction (str): One of "lower_is_better" or "higher_is_better".
        size: Font size for the arrow (default Pt(18)).
    
    Returns:
        The textbox shape containing the arrow.
    """
    indicator = get_status_indicator(status, direction)
    
    indicator_box = slide.shapes.add_textbox(left, top, Inches(0.3), Inches(0.3))
    indicator_frame = indicator_box.text_frame
    indicator_para = indicator_frame.paragraphs[0]
    indicator_para.text = indicator["arrow"]
    indicator_para.font.name = TITLE_FONT_NAME
    indicator_para.font.size = size
    indicator_para.font.bold = True
    indicator_para.font.color.rgb = indicator["color"]
    indicator_para.alignment = PP_ALIGN.CENTER
    
    return indicator_box


def create_metric_card_with_indicator(slide, left, top, width, height,
                                       value, label, metric_name,
                                       context=None, border_color=None,
                                       value_size=Pt(42), show_indicator=True):
    """Create a metric card with an optional status indicator arrow.
    
    This is an enhanced version of create_metric_card() that evaluates the
    metric against thresholds and displays an appropriate arrow indicator.
    
    Args:
        slide: The slide object to add the card to.
        left: Left position of the card.
        top: Top position of the card.
        width: Width of the card.
        height: Height of the card.
        value: The metric value to display (numeric).
        label: The label for the metric.
        metric_name (str): Key to look up in METRIC_THRESHOLDS.
        context: Optional context text below the label.
        border_color: Optional border color (defaults based on status or CS_BLUE).
        value_size: Font size for the value (default Pt(42)).
        show_indicator: Whether to show the trend indicator (default True).
    
    Returns:
        tuple: (card_shape, status_info) where status_info is the result from
               get_metric_status() or None if metric not found.
    """
    # Get metric status
    status_info = get_metric_status(metric_name, value) if show_indicator else None
    
    # Determine border color based on status if not explicitly provided
    if border_color is None:
        if status_info:
            indicator = get_status_indicator(status_info["status"], status_info["direction"])
            border_color = indicator["color"]
        else:
            border_color = CS_BLUE
    
    # Card background
    card_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = card_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    line = card_shape.line
    line.color.rgb = border_color
    line.width = Pt(3)
    
    # Calculate layout
    value_height = height * 0.45
    indicator_width = Inches(0.35) if status_info and show_indicator else 0
    
    # Value text (with room for indicator)
    value_box = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.15),
        width - Inches(0.2) - indicator_width, value_height
    )
    value_frame = value_box.text_frame
    value_para = value_frame.paragraphs[0]
    value_para.text = str(value)
    value_para.font.name = TITLE_FONT_NAME
    value_para.font.size = value_size
    value_para.font.bold = True
    value_para.font.color.rgb = CS_NAVY
    value_para.alignment = PP_ALIGN.CENTER
    
    # Add trend indicator if available
    if status_info and show_indicator:
        indicator_left = left + width - indicator_width - Inches(0.1)
        indicator_top = top + Inches(0.2)
        add_trend_indicator(
            slide, indicator_left, indicator_top,
            status_info["status"], status_info["direction"],
            size=Pt(24)
        )
    
    # Label
    label_top = top + value_height + Inches(0.1)
    label_box = slide.shapes.add_textbox(
        left + Inches(0.1), label_top,
        width - Inches(0.2), Inches(0.4)
    )
    label_frame = label_box.text_frame
    label_para = label_frame.paragraphs[0]
    label_para.text = label
    label_para.font.name = TITLE_FONT_NAME
    label_para.font.size = Pt(14)
    label_para.font.bold = True
    label_para.font.color.rgb = border_color
    label_para.alignment = PP_ALIGN.CENTER
    
    # Context (if provided)
    if context:
        context_top = label_top + Inches(0.35)
        context_box = slide.shapes.add_textbox(
            left + Inches(0.1), context_top,
            width - Inches(0.2), Inches(0.4)
        )
        context_frame = context_box.text_frame
        context_para = context_frame.paragraphs[0]
        context_para.text = context
        context_para.font.name = BODY_FONT_NAME
        context_para.font.size = Pt(11)
        context_para.font.color.rgb = CS_SLATE
        context_para.alignment = PP_ALIGN.CENTER
    
    return card_shape, status_info


def get_status_text(status, direction, metric_label=None):
    """Generate human-readable status text for narrative insight bars.
    
    Args:
        status (str): One of "good", "warning", or "bad".
        direction (str): One of "lower_is_better" or "higher_is_better".
        metric_label (str, optional): The metric name for more specific messaging.
    
    Returns:
        str: A brief status description like "above target" or "needs attention".
    """
    if status == "good":
        if direction == "higher_is_better":
            return "exceeding target"
        else:
            return "below target (optimal)"
    elif status == "warning":
        return "approaching threshold"
    else:  # bad
        if direction == "higher_is_better":
            return "below target"
        else:
            return "elevated (needs attention)"
